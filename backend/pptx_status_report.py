"""Générateur Status Report PPT — MARCEL."""
from pptx_base import *
from pptx_base import _brand, _set_brand, _blank_slide, _rect, _tb, _clear, _run, _footer, _CURRENT_BRAND

# ═══════════════════════════════════════════════════════════════════════════════
# STATUS REPORT PPT — 6 slides
# ═══════════════════════════════════════════════════════════════════════════════

WEATHER_META = {
    "soleil":  {"label": "Sous contrôle",     "char": "●", "color": RGBColor(0xFF, 0xD7, 0x00), "bg": RGBColor(0xFF, 0xFD, 0xE7)},
    "nuage":   {"label": "Point d'attention", "char": "●", "color": RGBColor(0x78, 0x91, 0x9C), "bg": RGBColor(0xF0, 0xF4, 0xF8)},
    "pluie":   {"label": "Problème avéré",    "char": "●", "color": RGBColor(0x29, 0x7D, 0xD5), "bg": RGBColor(0xEB, 0xF4, 0xFF)},
    "orage":   {"label": "Critique",          "char": "●", "color": RGBColor(0xEF, 0x44, 0x44), "bg": RGBColor(0xFF, 0xF0, 0xF0)},
    "gel":     {"label": "Bloqué",            "char": "●", "color": RGBColor(0x1E, 0x3A, 0x8A), "bg": RGBColor(0xE8, 0xEE, 0xF9)},
}

WEATHER_EMOJI = {
    "soleil": "SOLEIL",
    "nuage":  "NUAGE",
    "pluie":  "PLUIE",
    "orage":  "ORAGE",
    "gel":    "GEL",
}


def _sr_header(slide, title: str, brand: dict):
    """Bandeau titre slide Status Report."""
    _rect(slide, 0, 0, SW, Inches(0.65), fill=brand["primary"])
    tb = _tb(slide, Inches(0.4), Inches(0.07), SW - Inches(0.8), Inches(0.5))
    _clear(tb.text_frame)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.name = brand["font"]
    r.font.color.rgb = WHITE


def _sr_footer(slide, brand: dict, project_name: str):
    """Footer Status Report."""
    now_str = datetime.now().strftime("%d/%m/%Y")
    _footer(slide, f"{project_name} · Status Report · {now_str} · {brand['company_name']}")


def _weather_cell(slide, left, top, width, height, indicator_label: str, level: str, comment: str, brand: dict):
    """Dessine une cellule météo (1/4 de la grille 2×2)."""
    meta = WEATHER_META.get(level, WEATHER_META["gel"])

    # Fond de cellule
    _rect(slide, left, top, width, height, fill=meta["bg"], no_line=False, line_color=BORDER, line_pt=0.75)

    pad = Inches(0.12)
    inner_w = width - 2 * pad

    # Titre indicateur
    tb_title = _tb(slide, left + pad, top + Inches(0.1), inner_w, Inches(0.3))
    _clear(tb_title.text_frame)
    p = tb_title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = indicator_label
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.name = brand["font"]
    r.font.color.rgb = DARK

    # Icône météo (cercle coloré)
    circle_size = Inches(0.75)
    cx = left + (width - circle_size) / 2
    cy = top + Inches(0.42)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu
    circle = slide.shapes.add_shape(
        9,  # MSO_SHAPE_TYPE.OVAL
        cx, cy, circle_size, circle_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = meta["color"]
    circle.line.width = 0

    # Texte niveau
    tb_level = _tb(slide, left + pad, top + Inches(1.22), inner_w, Inches(0.3))
    _clear(tb_level.text_frame)
    p2 = tb_level.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = meta["label"].upper()
    r2.font.size = Pt(9)
    r2.font.bold = True
    r2.font.name = brand["font"]
    r2.font.color.rgb = meta["color"]

    # Commentaire CP
    if comment:
        tb_com = _tb(slide, left + pad, top + Inches(1.52), inner_w, Inches(0.65))
        _clear(tb_com.text_frame)
        p3 = tb_com.text_frame.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = trunc(comment, 120)
        r3.font.size = Pt(8)
        r3.font.italic = True
        r3.font.name = brand["font"]
        r3.font.color.rgb = MID


def _sr_table(slide, left, top, width, headers: list, rows: list, col_widths: list, brand: dict):
    """Tableau générique pour slides 3-6."""
    row_h = Inches(0.28)
    hdr_h = Inches(0.32)
    total_h = hdr_h + row_h * min(len(rows), 16)

    # En-têtes
    x = left
    _rect(slide, left, top, width, hdr_h, fill=brand["primary"])
    for i, (hdr, cw) in enumerate(zip(headers, col_widths)):
        tb = _tb(slide, x + Inches(0.06), top + Inches(0.04), cw - Inches(0.12), hdr_h)
        _clear(tb.text_frame)
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = hdr
        r.font.size = Pt(7.5)
        r.font.bold = True
        r.font.name = brand["font"]
        r.font.color.rgb = WHITE
        x += cw

    # Lignes
    for ri, row in enumerate(rows[:16]):
        y = top + hdr_h + ri * row_h
        bg = BG if ri % 2 == 0 else WHITE
        _rect(slide, left, y, width, row_h, fill=bg, no_line=False, line_color=BORDER, line_pt=0.25)
        x = left
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            val = str(cell) if cell is not None else "—"
            tb = _tb(slide, x + Inches(0.04), y + Inches(0.03), cw - Inches(0.08), row_h)
            _clear(tb.text_frame)
            p = tb.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = trunc(val, 45)
            r.font.size = Pt(7)
            r.font.name = brand["font"]
            r.font.color.rgb = DARK
            x += cw


# ── Slide 1 — Garde ──────────────────────────────────────────────────────────
def _sr_slide1_cover(prs, project: dict, program_name: str, cp_name: str, brand: dict) -> None:
    slide = _blank_slide(prs)
    # Fond plein primary
    _rect(slide, 0, 0, SW, SH, fill=brand["primary"])
    # Bandeau blanc bas
    _rect(slide, 0, SH - Inches(2.0), SW, Inches(2.0), fill=WHITE)

    # Logo tenant
    logo_b64 = brand.get("logo_base64")
    if logo_b64:
        try:
            logo_data = base64.b64decode(logo_b64)
            logo_buf = io.BytesIO(logo_data)
            logo_w = Inches(1.5)
            slide.shapes.add_picture(logo_buf, Inches(0.5), Inches(0.4), width=logo_w)
        except Exception:
            pass

    # Titre projet
    tb1 = _tb(slide, Inches(1.0), Inches(1.8), SW - Inches(2.0), Inches(1.2))
    _clear(tb1.text_frame)
    p1 = tb1.text_frame.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "STATUS REPORT"
    r1.font.size = Pt(18)
    r1.font.bold = True
    r1.font.name = brand["font"]
    r1.font.color.rgb = WHITE

    tb2 = _tb(slide, Inches(1.0), Inches(2.8), SW - Inches(2.0), Inches(1.0))
    _clear(tb2.text_frame)
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = project.get("name", "Projet")
    r2.font.size = Pt(28)
    r2.font.bold = True
    r2.font.name = brand["font"]
    r2.font.color.rgb = WHITE

    # Infos bas
    date_str = datetime.now().strftime("%d %B %Y").capitalize()
    infos = [
        ("Programme", program_name),
        ("Chef de Projet", cp_name),
        ("Date du report", date_str),
        ("Statut projet", project.get("status", "—")),
    ]
    y = SH - Inches(1.75)
    for label, val in infos:
        tb = _tb(slide, Inches(1.5), y, SW - Inches(3.0), Inches(0.32))
        _clear(tb.text_frame)
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{label} : {val}"
        r.font.size = Pt(10)
        r.font.name = brand["font"]
        r.font.color.rgb = DARK
        y += Inches(0.36)


# ── Slide 2 — Météo 2×2 ───────────────────────────────────────────────────────
def _sr_slide2_meteo(prs, weather: dict, brand: dict) -> None:
    slide = _blank_slide(prs)
    _sr_header(slide, "Météo Projet", brand)

    margin = Inches(0.35)
    available_w = SW - 2 * margin
    available_h = SH - Inches(0.65) - Inches(0.35) - margin
    cell_w = available_w / 2
    cell_h = available_h / 2

    indicators = [
        ("Périmètre",           weather.get("perimeter", {})),
        ("Budget",              weather.get("budget", {})),
        ("Calendrier",          weather.get("calendar", {})),
        ("Changement de scope", weather.get("scope_change", {})),
    ]

    for i, (label, info) in enumerate(indicators):
        col = i % 2
        row = i // 2
        left = margin + col * cell_w
        top = Inches(0.65) + Inches(0.1) + row * cell_h
        level = info.get("level", "gel")
        comment = info.get("comment", "")
        _weather_cell(slide, left, top, cell_w - Inches(0.05), cell_h - Inches(0.05),
                      label, level, comment, brand)

    _sr_footer(slide, brand, "")


# ── Slide 3 — Jalons livrés ───────────────────────────────────────────────────
def _sr_slide3_milestones_done(prs, milestones: list, brand: dict) -> None:
    slide = _blank_slide(prs)
    _sr_header(slide, "Jalons livrés", brand)

    today = datetime.now().date()
    done_ms = []
    for m in milestones:
        if m.get("status") != "done":
            continue
        date_str = m.get("date_reel") or m.get("date_forecast") or ""
        try:
            d = datetime.fromisoformat(date_str[:10]).date()
        except Exception:
            d = None
        baseline_str = m.get("date_baseline") or ""
        try:
            b = datetime.fromisoformat(baseline_str[:10]).date()
            ecart = (d - b).days if d and b else None
        except Exception:
            ecart = None
        done_ms.append({
            "name": m.get("name", ""),
            "family": m.get("family", ""),
            "date_baseline": fmt_date(baseline_str),
            "date_reel": fmt_date(date_str),
            "ecart": f"{ecart:+d} j" if ecart is not None else "—",
            "_sort": d or today,
        })

    done_ms.sort(key=lambda x: x["_sort"], reverse=True)

    if not done_ms:
        tb = _tb(slide, Inches(0.5), Inches(1.0), SW - Inches(1.0), Inches(0.5))
        _clear(tb.text_frame)
        _run(tb.text_frame, "Aucun jalon livré.", 10, italic=True, color=MID)
        return

    headers = ["Jalon", "Famille", "Date prévue", "Date réelle", "Écart"]
    col_widths = [Inches(4.5), Inches(1.5), Inches(1.8), Inches(1.8), Inches(1.0)]
    rows = [[m["name"], m["family"], m["date_baseline"], m["date_reel"], m["ecart"]] for m in done_ms]
    _sr_table(slide, Inches(0.35), Inches(0.75), SW - Inches(0.7), headers, rows, col_widths, brand)
    _sr_footer(slide, brand, "")


# ── Slide 4 — Jalons futurs ───────────────────────────────────────────────────
def _sr_slide4_milestones_future(prs, milestones: list, res_map: dict, brand: dict) -> None:
    slide = _blank_slide(prs)
    _sr_header(slide, "Jalons à venir", brand)

    today = datetime.now().date()
    future_ms = []
    for m in milestones:
        if m.get("status") == "done":
            continue
        date_str = m.get("date_forecast") or m.get("date_baseline") or ""
        try:
            d = datetime.fromisoformat(date_str[:10]).date()
        except Exception:
            continue
        if d >= today:
            future_ms.append({
                "name": m.get("name", ""),
                "family": m.get("family", ""),
                "date_forecast": fmt_date(date_str),
                "owner": res_map.get(m.get("owner_id", ""), "—"),
                "bloquant": "Oui" if m.get("is_blocking") else "Non",
                "_sort": d,
            })

    future_ms.sort(key=lambda x: x["_sort"])

    if not future_ms:
        tb = _tb(slide, Inches(0.5), Inches(1.0), SW - Inches(1.0), Inches(0.5))
        _clear(tb.text_frame)
        _run(tb.text_frame, "Aucun jalon à venir.", 10, italic=True, color=MID)
        return

    headers = ["Jalon", "Famille", "Date prévue", "Owner", "Bloquant"]
    col_widths = [Inches(4.5), Inches(1.5), Inches(1.8), Inches(2.5), Inches(1.0)]
    rows = [[m["name"], m["family"], m["date_forecast"], m["owner"], m["bloquant"]] for m in future_ms]
    _sr_table(slide, Inches(0.35), Inches(0.75), SW - Inches(0.7), headers, rows, col_widths, brand)
    _sr_footer(slide, brand, "")


# ── Slide 5 — Jalons métiers ──────────────────────────────────────────────────
def _sr_slide5_milestones_business(prs, milestones: list, brand: dict) -> None:
    slide = _blank_slide(prs)
    _sr_header(slide, "Jalons métiers & transverses", brand)

    business_ms = []
    for m in milestones:
        is_transversal = m.get("family") in ("transversal", "regulatory", "dependency", "decomm")
        is_key = m.get("attribute") in ("critical", "strategic")
        if is_transversal or is_key:
            date_str = m.get("date_forecast") or m.get("date_baseline") or ""
            business_ms.append({
                "name": m.get("name", ""),
                "family": m.get("family", ""),
                "attribute": m.get("attribute", "—"),
                "date": fmt_date(date_str),
                "status": m.get("status", ""),
            })

    if not business_ms:
        tb = _tb(slide, Inches(0.5), Inches(1.0), SW - Inches(1.0), Inches(0.5))
        _clear(tb.text_frame)
        _run(tb.text_frame, "Aucun jalon métier / transverse défini.", 10, italic=True, color=MID)
        return

    headers = ["Jalon", "Famille", "Attribut", "Date", "Statut"]
    col_widths = [Inches(4.5), Inches(1.8), Inches(1.5), Inches(1.8), Inches(1.2)]
    rows = [[m["name"], m["family"], m["attribute"], m["date"], m["status"]] for m in business_ms]
    _sr_table(slide, Inches(0.35), Inches(0.75), SW - Inches(0.7), headers, rows, col_widths, brand)
    _sr_footer(slide, brand, "")


# ── Slide 6 — Risques ────────────────────────────────────────────────────────
def _sr_slide6_risks(prs, risks: list, res_map: dict, brand: dict) -> None:
    slide = _blank_slide(prs)
    _sr_header(slide, "Risques ouverts", brand)

    open_risks = [r for r in risks if r.get("status") not in ("clos", "closed")]
    for r in open_risks:
        p = r.get("probability", 0) or 0
        i = r.get("impact", 0) or 0
        r["_crit"] = p * i

    open_risks.sort(key=lambda x: x["_crit"], reverse=True)

    if not open_risks:
        tb = _tb(slide, Inches(0.5), Inches(1.0), SW - Inches(1.0), Inches(0.5))
        _clear(tb.text_frame)
        _run(tb.text_frame, "Aucun risque ouvert.", 10, italic=True, color=MID)
        return

    headers = ["Risque", "P", "I", "P×I", "Catégorie", "Owner", "Mitigation"]
    col_widths = [Inches(3.0), Inches(0.5), Inches(0.5), Inches(0.6), Inches(1.5), Inches(2.0), Inches(3.0)]
    rows = [
        [
            r.get("name", r.get("title", "")),
            str(r.get("probability", "—")),
            str(r.get("impact", "—")),
            str(r["_crit"]) if r["_crit"] else "—",
            r.get("category", "—"),
            res_map.get(r.get("owner_id", ""), r.get("owner", "—")),
            r.get("mitigation_plan", r.get("mitigation", "—")),
        ]
        for r in open_risks
    ]
    _sr_table(slide, Inches(0.2), Inches(0.75), SW - Inches(0.4), headers, rows, col_widths, brand)
    _sr_footer(slide, brand, "")


# ── Fonction principale Status Report ────────────────────────────────────────
def generate_status_report_pptx(
    project: dict,
    program_name: str,
    cp_name: str,
    weather: dict,
    milestones: list,
    risks: list,
    res_map: dict,
    branding: dict | None,
) -> bytes:
    """Génère le PPT Status Report (6 slides) et retourne les bytes."""
    brand = _brand(branding)
    _set_brand(brand)  # Met à jour pptx_base._CURRENT_BRAND

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    _sr_slide1_cover(prs, project, program_name, cp_name, brand)
    _sr_slide2_meteo(prs, weather, brand)
    _sr_slide3_milestones_done(prs, milestones, brand)
    _sr_slide4_milestones_future(prs, milestones, res_map, brand)
    _sr_slide5_milestones_business(prs, milestones, brand)
    _sr_slide6_risks(prs, risks, res_map, brand)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
