"""Générateur PowerPoint charte MARCEL (16:9)."""
import io
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

INDIGO = RGBColor(0x35, 0x2C, 0x6E)
BLUE = RGBColor(0x2E, 0x5F, 0xE8)
BG = RGBColor(0xF7, 0xF6, 0xFB)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x26, 0x24, 0x3A)
MUTED = RGBColor(0x8A, 0x87, 0xA0)
BORDER = RGBColor(0xE8, 0xE6, 0xF0)
GREEN = RGBColor(0x3F, 0x8A, 0x34)
AMBER = RGBColor(0xB7, 0x79, 0x1F)
RED = RGBColor(0xCC, 0x4F, 0x45)
FONT = "Outfit"

SW, SH = Inches(13.333), Inches(7.5)


def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _round_rect(slide, x, y, w, h, color, radius=0.08):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    _set_fill(sh, color)
    sh.shadow.inherit = False
    return sh


def _text(slide, x, y, w, h, text, size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
          font=FONT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font
    return tb


class MarcelDeck:
    def __init__(self, tenant_name: str = ""):
        self.prs = Presentation()
        self.prs.slide_width = SW
        self.prs.slide_height = SH
        self.blank = self.prs.slide_layouts[6]
        self.tenant = tenant_name
        self.page = 0

    def _new(self, bg=BG):
        slide = self.prs.slides.add_slide(self.blank)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
        _set_fill(rect, bg)
        rect.shadow.inherit = False
        return slide

    def _footer(self, slide):
        self.page += 1
        _text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.35),
              f"MARCEL · {self.tenant} · {date.today().strftime('%d/%m/%Y')}",
              size=9, color=MUTED)
        _text(slide, Inches(12.2), Inches(7.05), Inches(0.6), Inches(0.35),
              str(self.page), size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    def _header(self, slide, title, subtitle=None):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55), Inches(0.09), Inches(0.62))
        _set_fill(bar, BLUE)
        bar.shadow.inherit = False
        _text(slide, Inches(0.85), Inches(0.48), Inches(11), Inches(0.5), title, size=24, bold=True, color=INDIGO)
        if subtitle:
            _text(slide, Inches(0.85), Inches(0.98), Inches(11), Inches(0.35), subtitle, size=12, color=MUTED)

    # ── Slides types ──────────────────────────────────────────────────────────

    def cover(self, title: str, subtitle: str = "", meta: str = ""):
        slide = self._new(bg=INDIGO)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.9), SW, Inches(0.6))
        _set_fill(band, BLUE)
        band.shadow.inherit = False
        _text(slide, Inches(0.8), Inches(0.7), Inches(4), Inches(0.5), "MARCEL",
              size=18, bold=True, color=RGBColor(0xB7, 0xAE, 0xF7))
        _text(slide, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.6), title, size=44, bold=True,
              color=RGBColor(0xFF, 0xFF, 0xFF))
        if subtitle:
            _text(slide, Inches(0.8), Inches(4.1), Inches(11.5), Inches(0.6), subtitle, size=18,
                  color=RGBColor(0xC5, 0xBD, 0xF0))
        _text(slide, Inches(0.8), Inches(6.1), Inches(11), Inches(0.4),
              meta or f"{self.tenant} — {date.today().strftime('%d %B %Y')}",
              size=13, color=RGBColor(0x98, 0x95, 0xB0))
        return slide

    def section(self, title: str, subtitle: str = ""):
        slide = self._new()
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH)
        _set_fill(band, INDIGO)
        band.shadow.inherit = False
        _text(slide, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.0), title, size=34, bold=True, color=INDIGO)
        if subtitle:
            _text(slide, Inches(0.9), Inches(4.0), Inches(11), Inches(0.6), subtitle, size=15, color=MUTED)
        self._footer(slide)
        return slide

    def kpis(self, title: str, items: list, subtitle: str = None):
        """items: [{label, value, color?, hint?}] — max 5 par ligne."""
        slide = self._new()
        self._header(slide, title, subtitle)
        n = min(len(items), 5)
        gap = Inches(0.25)
        card_w = Emu(int((SW - Inches(1.2) - gap * (n - 1)) / n))
        x = Inches(0.6)
        for it in items[:5]:
            _round_rect(slide, x, Inches(1.7), card_w, Inches(1.9), CARD)
            _text(slide, x + Inches(0.25), Inches(1.95), card_w - Inches(0.5), Inches(0.35),
                  str(it["label"]).upper(), size=10, bold=True, color=MUTED)
            _text(slide, x + Inches(0.25), Inches(2.35), card_w - Inches(0.5), Inches(0.75),
                  str(it["value"]), size=30, bold=True, color=it.get("color", INDIGO))
            if it.get("hint"):
                _text(slide, x + Inches(0.25), Inches(3.15), card_w - Inches(0.5), Inches(0.35),
                      str(it["hint"]), size=10, color=MUTED)
            x = Emu(int(x + card_w + gap))
        self._footer(slide)
        return slide

    def table(self, title: str, headers: list, rows: list, subtitle: str = None,
              col_widths: list = None, cell_colors: dict = None, start_y: float = 1.6):
        """rows: list[list[str]] — cell_colors: {(row_idx, col_idx): RGBColor}."""
        slide = self._new()
        self._header(slide, title, subtitle)
        n_rows, n_cols = len(rows) + 1, len(headers)
        total_w = SW - Inches(1.2)
        table_h = Inches(min(0.42 * n_rows + 0.1, 5.1))
        shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(start_y), total_w, table_h)
        tbl = shape.table
        if col_widths:
            s = sum(col_widths)
            for i, cw in enumerate(col_widths):
                tbl.columns[i].width = Emu(int(total_w * cw / s))
        # Style d'en-tête
        for c, h in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = INDIGO
            cell.margin_top = cell.margin_bottom = Pt(4)
            tf = cell.text_frame
            tf.paragraphs[0].text = str(h)
            for p in tf.paragraphs:
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.name = FONT
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD if r % 2 else BG
                cell.margin_top = cell.margin_bottom = Pt(3)
                tf = cell.text_frame
                tf.paragraphs[0].text = str(val)
                for p in tf.paragraphs:
                    p.font.size = Pt(10.5)
                    p.font.name = FONT
                    p.font.color.rgb = (cell_colors or {}).get((r - 1, c), TEXT)
        self._footer(slide)
        return slide

    def bullets(self, title: str, lines: list, subtitle: str = None):
        slide = self._new()
        self._header(slide, title, subtitle)
        card = _round_rect(slide, Inches(0.6), Inches(1.7), SW - Inches(1.2), Inches(4.9), CARD)
        tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), SW - Inches(2.0), Inches(4.3))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            r = p.add_run()
            r.text = f"•  {line}"
            r.font.size = Pt(14)
            r.font.name = FONT
            r.font.color.rgb = TEXT
        self._footer(slide)
        return slide

    def to_bytes(self) -> bytes:
        buf = io.BytesIO()
        self.prs.save(buf)
        return buf.getvalue()
