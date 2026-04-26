"""Tests for COPIL PPTX export endpoint"""
import pytest
import requests
import os
import io
from pptx import Presentation

BASE_URL = os.environ.get("REACT_APP_BACKEND_URL", "").rstrip("/")


@pytest.fixture(scope="module")
def auth_token():
    res = requests.post(f"{BASE_URL}/api/auth/login", json={"email": "admin@altair.fr", "password": "Admin2026!"})
    assert res.status_code == 200, f"Login failed: {res.text}"
    data = res.json()
    return data.get("access_token") or data.get("token")


@pytest.fixture(scope="module")
def project_ids(auth_token):
    res = requests.get(f"{BASE_URL}/api/projects", headers={"Authorization": f"Bearer {auth_token}"})
    assert res.status_code == 200
    projects = res.json()
    assert len(projects) >= 2, "Need at least 2 projects"
    return [p["project_id"] for p in projects[:3]]


@pytest.fixture(scope="module")
def pptx_response(auth_token, project_ids):
    """Shared PPTX response for multiple tests"""
    payload = {
        "project_ids": project_ids,
        "instance_name": "Test COPIL",
        "instance_date": "2026-04-26",
    }
    res = requests.post(
        f"{BASE_URL}/api/export/copil",
        json=payload,
        headers={"Authorization": f"Bearer {auth_token}"},
    )
    return res


# --- Backend tests ---

def test_export_copil_http_200(pptx_response):
    """POST /api/export/copil returns HTTP 200"""
    assert pptx_response.status_code == 200, f"Expected 200, got {pptx_response.status_code}: {pptx_response.text[:200]}"


def test_export_copil_content_type(pptx_response):
    """Response content-type is PPTX"""
    ct = pptx_response.headers.get("content-type", "")
    assert "presentationml" in ct or "application/octet-stream" in ct, f"Unexpected content-type: {ct}"


def test_export_copil_content_disposition_format(pptx_response):
    """Content-Disposition filename follows COPIL_[date]_[slug].pptx format"""
    cd = pptx_response.headers.get("content-disposition", "")
    print(f"Content-Disposition: {cd}")
    assert "COPIL_" in cd, f"Missing COPIL_ prefix in: {cd}"
    assert ".pptx" in cd, f"Missing .pptx extension in: {cd}"
    # Date should come before slug: COPIL_2026-04-26_...
    assert "2026-04-26" in cd, f"Date not found in filename: {cd}"


def test_export_copil_binary_size(pptx_response):
    """Response body is a non-empty PPTX binary (>10KB)"""
    assert len(pptx_response.content) > 10000, f"PPTX too small: {len(pptx_response.content)} bytes"


def test_export_copil_at_least_7_slides(pptx_response):
    """PPTX has >= 7 slides"""
    prs = Presentation(io.BytesIO(pptx_response.content))
    n = len(prs.slides)
    print(f"Slide count: {n}")
    assert n >= 7, f"Expected >= 7 slides, got {n}"


def test_export_copil_last_slide_is_cloture(pptx_response):
    """Last slide contains 'Points d'attention CIO' as title"""
    prs = Presentation(io.BytesIO(pptx_response.content))
    last_slide = prs.slides[-1]
    all_text = ""
    for shape in last_slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    all_text += run.text + " "
    print(f"Last slide text: {all_text[:300]}")
    assert "Points d'attention CIO" in all_text or "attention CIO" in all_text, \
        f"'Points d'attention CIO' not found in last slide. Text: {all_text[:300]}"


def test_export_copil_with_governance_filter(auth_token):
    """When governance_id provided, endpoint returns HTTP 200"""
    gov_res = requests.get(f"{BASE_URL}/api/governance", headers={"Authorization": f"Bearer {auth_token}"})
    if gov_res.status_code != 200 or not gov_res.json():
        pytest.skip("No governance instances available")
    
    gov_data = gov_res.json()[0]
    governance_id = gov_data.get("governance_id") or gov_data.get("id")
    
    proj_res = requests.get(f"{BASE_URL}/api/projects", headers={"Authorization": f"Bearer {auth_token}"})
    project_ids = [p["project_id"] for p in proj_res.json()[:2]]
    
    payload = {
        "project_ids": project_ids,
        "instance_name": "Test Governance COPIL",
        "instance_date": "2026-04-26",
        "governance_id": governance_id,
    }
    pptx_res = requests.post(
        f"{BASE_URL}/api/export/copil",
        json=payload,
        headers={"Authorization": f"Bearer {auth_token}"},
    )
    assert pptx_res.status_code == 200
    assert len(pptx_res.content) > 5000


def test_export_copil_unauthenticated():
    """Unauthenticated requests return 401 or 403"""
    res = requests.post(f"{BASE_URL}/api/export/copil", json={
        "project_ids": ["fake-id"],
        "instance_name": "Test",
        "instance_date": "2026-04-26"
    })
    assert res.status_code in (401, 403), f"Expected 401/403, got {res.status_code}"
