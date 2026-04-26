"""
Tests: PPT Branding, Timesheet 2-step workflow, Admin config, MilestoneModal merge
Iteration 29 — Projetenne Admin Back-office integration points
"""
import pytest
import requests
import os
import io

BASE_URL = os.environ.get("REACT_APP_BACKEND_URL", "http://localhost:8001").rstrip("/")


def login(email, password):
    r = requests.post(f"{BASE_URL}/api/auth/login", json={"email": email, "password": password})
    assert r.status_code == 200, f"Login failed {r.status_code}: {r.text}"
    return r.json()["access_token"]


@pytest.fixture(scope="module")
def admin_token():
    return login("admin@altair.fr", "Admin2026!")


@pytest.fixture(scope="module")
def pmo_token():
    return login("cp@altair.fr", "Altair2026!")


# ─── Test 1: Admin config retourne les bons champs ───────────────────────────

class TestAdminConfig:
    """Admin config sanity checks"""

    def test_config_returns_200(self, admin_token):
        r = requests.get(
            f"{BASE_URL}/api/admin/config",
            headers={"Authorization": f"Bearer {admin_token}"},
        )
        assert r.status_code == 200, f"Expected 200, got {r.status_code}: {r.text}"

    def test_ppt_branding_company_name(self, admin_token):
        r = requests.get(
            f"{BASE_URL}/api/admin/config",
            headers={"Authorization": f"Bearer {admin_token}"},
        )
        data = r.json()
        company_name = (
            data.get("ppt_branding", {}).get("company_name")
            or data.get("settings", {}).get("ppt_branding", {}).get("company_name")
        )
        assert company_name == "Groupe Altair Industries", (
            f"Expected 'Groupe Altair Industries', got '{company_name}'. Full response: {data}"
        )

    def test_timesheet_validation_steps_2(self, admin_token):
        r = requests.get(
            f"{BASE_URL}/api/admin/config",
            headers={"Authorization": f"Bearer {admin_token}"},
        )
        data = r.json()
        # Try both possible response structures
        steps = (
            data.get("workflows", {}).get("timesheet", {}).get("validation_steps")
            or data.get("settings", {}).get("workflows", {}).get("timesheet", {}).get("validation_steps")
        )
        assert steps == 2, (
            f"Expected validation_steps=2, got {steps}. Full response keys: {list(data.keys())}"
        )


# ─── Test 2: Export PPTX — branding ──────────────────────────────────────────

class TestPPTXBranding:
    """PPTX export branding tests"""

    def _get_first_project(self, admin_token):
        r = requests.get(
            f"{BASE_URL}/api/projects",
            headers={"Authorization": f"Bearer {admin_token}"},
        )
        assert r.status_code == 200
        data = r.json()
        projects = data if isinstance(data, list) else data.get("items", [])
        assert len(projects) > 0, "No projects found"
        return [projects[0]["project_id"]]

    def test_export_copil_returns_pptx(self, admin_token):
        project_ids = self._get_first_project(admin_token)
        r = requests.post(
            f"{BASE_URL}/api/export/copil",
            headers={"Authorization": f"Bearer {admin_token}"},
            json={
                "project_ids": project_ids,
                "instance_name": "COPIL Test Branding",
                "instance_date": "2026-02-01",
            },
        )
        assert r.status_code == 200, f"Export failed {r.status_code}: {r.text[:300]}"
        ct = r.headers.get("content-type", "")
        assert "openxml" in ct or "pptx" in ct or "octet" in ct, f"Unexpected content-type: {ct}"

    def test_export_pptx_contains_company_name(self, admin_token):
        """Verify PPTX contains 'Groupe Altair Industries' in text runs."""
        project_ids = self._get_first_project(admin_token)
        r = requests.post(
            f"{BASE_URL}/api/export/copil",
            headers={"Authorization": f"Bearer {admin_token}"},
            json={
                "project_ids": project_ids,
                "instance_name": "COPIL Branding Check",
                "instance_date": "2026-02-01",
            },
        )
        assert r.status_code == 200
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            pytest.skip("python-pptx not available")

        prs = Presentation(io.BytesIO(r.content))
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            all_text.append(run.text)

        full = " ".join(all_text)
        assert "Groupe Altair Industries" in full, (
            f"'Groupe Altair Industries' not found in PPTX text. Sample text: {full[:500]}"
        )

    def test_export_pptx_primary_color_in_header(self, admin_token):
        """Verify primary color #0B2545 is used in slide shapes (header backgrounds)."""
        project_ids = self._get_first_project(admin_token)
        r = requests.post(
            f"{BASE_URL}/api/export/copil",
            headers={"Authorization": f"Bearer {admin_token}"},
            json={
                "project_ids": project_ids,
                "instance_name": "COPIL Color Check",
                "instance_date": "2026-02-01",
            },
        )
        assert r.status_code == 200
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
        except ImportError:
            pytest.skip("python-pptx not available")

        prs = Presentation(io.BytesIO(r.content))
        target = (0x0B, 0x25, 0x45)
        found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    fill = shape.fill
                    if fill.type is not None:
                        fg = fill.fore_color
                        if fg and fg.rgb == RGBColor(*target):
                            found = True
                            break
                except Exception:
                    pass
            if found:
                break
        assert found, "Primary color #0B2545 not found in any shape fill in PPTX"


# ─── Test 3: Timesheet workflow 2-étapes ─────────────────────────────────────

class TestTimesheetWorkflow2Steps:
    """PMO validation should go directly to validated (bypass CP) with 2-step config."""

    def test_pmo_validate_timesheets_flow(self, pmo_token):
        """PMO validation: submitted → validated (not cp_reviewed)."""
        # Get list of submitted timesheets
        r = requests.get(
            f"{BASE_URL}/api/timesheets/validation?view=pmo",
            headers={"Authorization": f"Bearer {pmo_token}"},
        )
        assert r.status_code == 200, f"Validation view failed: {r.status_code}: {r.text}"
        groups = r.json()

        submitted_groups = [g for g in groups if g.get("status") == "submitted"]
        if not submitted_groups:
            pytest.skip("No submitted timesheets available for PMO validation test")

        # Pick first submitted group's ts_ids
        ts_ids = submitted_groups[0]["ts_ids"]
        assert len(ts_ids) > 0

        r2 = requests.post(
            f"{BASE_URL}/api/timesheets/validate",
            headers={"Authorization": f"Bearer {pmo_token}"},
            json={"timesheet_ids": ts_ids},
        )
        assert r2.status_code == 200, f"Validate failed: {r2.status_code}: {r2.text}"
        result = r2.json()
        # PMO should go directly to validated
        assert result.get("validated", 0) > 0, (
            f"Expected validated > 0, got: {result}"
        )
        assert result.get("advanced_to_cp_reviewed", 0) == 0, (
            f"PMO should NOT advance to cp_reviewed (2-step), got: {result}"
        )

    def test_validation_steps_from_config(self, pmo_token):
        """Config endpoint returns validation_steps=2."""
        # Using pmo_token to check via admin config (pmo may not have access - try admin)
        r = requests.get(
            f"{BASE_URL}/api/admin/config",
            headers={"Authorization": f"Bearer {pmo_token}"},
        )
        # PMO may get 403, that's ok — main check is done via admin_token in TestAdminConfig
        if r.status_code == 403:
            pytest.skip("PMO user does not have access to /api/admin/config")
        assert r.status_code == 200
        data = r.json()
        steps = (
            data.get("workflows", {}).get("timesheet", {}).get("validation_steps")
            or data.get("settings", {}).get("workflows", {}).get("timesheet", {}).get("validation_steps")
        )
        assert steps == 2


# ─── Test 4: Backend import sanity ───────────────────────────────────────────

class TestBackendSanity:
    """Backend is running, timesheets service has no import errors."""

    def test_health(self):
        r = requests.get(f"{BASE_URL}/api/health")
        assert r.status_code == 200

    def test_timesheets_endpoint_accessible(self, admin_token):
        # Simple check that the timesheets module loaded correctly
        r = requests.get(
            f"{BASE_URL}/api/timesheets/validation?view=pmo",
            headers={"Authorization": f"Bearer {admin_token}"},
        )
        assert r.status_code in (200, 404), f"Unexpected: {r.status_code}: {r.text}"
