"""Générateur PowerPoint COPIL — MARCEL (v2)"""
import io
import base64
from datetime import datetime
from typing import List, Optional

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---- Slide dimensions (widescreen 16:9) ----
SW = Inches(13.33)
SH = Inches(7.5)

# ---- Palette — fond blanc, accents Navy #0B2545, RAG couleurs corrigées ----
NAVY        = RGBColor(0x0B, 0x25, 0x45)
BLUE        = RGBColor(0x00, 0x52, 0xCC)
LIGHT_BLUE  = RGBColor(0xEB, 0xF2, 0xFF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG          = RGBColor(0xF8, 0xF9, 0xFA)
DARK        = RGBColor(0x1E, 0x29, 0x3B)
MID         = RGBColor(0x64, 0x74, 0x8B)
LIGHT       = RGBColor(0x94, 0xA3, 0xB8)
BORDER      = RGBColor(0xE2, 0xE8, 0xF0)
RED         = RGBColor(0xEF, 0x44, 0x44)
ORANGE_C    = RGBColor(0xF5, 0x9E, 0x0B)
GREEN_C     = RGBColor(0x10, 0xB9, 0x81)
LIGHT_RED   = RGBColor(0xFE, 0xF2, 0xF2)
LIGHT_AMBER = RGBColor(0xFF, 0xF7, 0xED)
LIGHT_GREEN = RGBColor(0xF0, 0xFD, 0xF4)

DECISION_STATUS_COLORS = {
    "proposée":  RGBColor(0xE6, 0xF0, 0xFF),
    "prise":     RGBColor(0xED, 0xE9, 0xFE),
    "en_cours":  RGBColor(0xFE, 0xF3, 0xC7),
    "appliquée": RGBColor(0xDC, 0xFC, 0xE7),
    "reportée":  RGBColor(0xF1, 0xF5, 0xF9),
    "annulée":   RGBColor(0xFE, 0xE2, 0xE2),
}

FONT = "Arial"

# ── Branding dynamique (mis à jour à chaque génération) ──────────────────────
_CURRENT_BRAND: dict = {}


# ---- Branding helpers ----

def _parse_rgb(hex_str, default: RGBColor) -> RGBColor:
    """Parse une couleur hex (#RRGGBB) en RGBColor avec fallback."""
    if not hex_str:
        return default
    try:
        h = str(hex_str).lstrip('#')
        if len(h) == 6:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        pass
    return default


def _brand(branding: dict | None) -> dict:
    """Construit le dictionnaire de couleurs de marque depuis la config tenant."""
    b = branding or {}
    return {
        "primary":      _parse_rgb(b.get("primary_color"), NAVY),
        "secondary":    _parse_rgb(b.get("secondary_color"), BLUE),
        "accent":       _parse_rgb(b.get("accent_color"), GREEN_C),
        "company_name": b.get("company_name") or "MARCEL",
        "font":         b.get("font") or FONT,
        "logo_base64":  b.get("logo_base64"),
    }



def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, left, top, width, height, fill=None, no_line=True, line_color=None, line_pt=0.5):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if no_line:
        shape.line.width = 0
    elif line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    return shape


def _tb(slide, left, top, width, height, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = wrap
    return tb


def _clear(tf):
    for p in tf.paragraphs:
        p.clear()


def _run(tf, text, size=9, bold=False, color=None, align=PP_ALIGN.LEFT, space_before=0, italic=False):
    """Add a paragraph with a single run to a text frame (brand font)."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = _CURRENT_BRAND.get("font", FONT) if _CURRENT_BRAND else FONT
    r.font.color.rgb = color or DARK
    return p


def keur(val):
    if val is None:
        return "—"
    return f"{int(val / 1000):,}".replace(",", "\u202f") + " K€"


def eur(val):
    """Format as xxx K€ or xxx € depending on magnitude."""
    if val is None:
        return "—"
    if abs(val) >= 1000:
        return keur(val)
    return f"{int(val):,}".replace(",", "\u202f") + " €"


def jh_fmt(val):
    if val is None:
        return "—"
    return f"{val:.1f} JH" if val != int(val) else f"{int(val)} JH"


def fmt_date(d):
    if not d:
        return "—"
    try:
        return datetime.fromisoformat(str(d).replace("Z", "")).strftime("%d/%m/%Y")
    except Exception:
        return str(d)[:10] if d else "—"


def trunc(s, n=55):
    if not s:
        return "—"
    return s[:n] + "…" if len(s) > n else s


def crit_color(c):
    return RED if c >= 16 else ORANGE_C if c >= 7 else GREEN_C


def crit_bg(c):
    return LIGHT_RED if c >= 16 else LIGHT_AMBER if c >= 7 else LIGHT_GREEN


def rag_color(rag):
    return {"green": GREEN_C, "orange": ORANGE_C, "red": RED}.get(str(rag), MID)


def rag_label(rag):
    return {"green": "Vert", "orange": "Orange", "red": "Rouge"}.get(str(rag), "—")


def status_label(s):
    return {
        "en_preparation": "En préparation", "actif": "Actif",
        "en_pause": "En pause", "cloture": "Clôturé", "archive": "Archivé",
    }.get(str(s), str(s) if s else "—")


def decision_status_label(s):
    return {
        "proposée": "Proposée", "prise": "Prise", "en_cours": "En cours",
        "appliquée": "Appliquée", "reportée": "Reportée", "annulée": "Annulée",
    }.get(str(s), str(s) if s else "—")


# ---- S2-04 — Slide Roadmap matplotlib ----

def add_slide_roadmap(prs, projects, all_milestones, instance_name, instance_date, dependencies=None, brand=None):
    """Génère un slide Roadmap avec un diagramme Gantt matplotlib multi-projets.
    Losanges colorés par famille (or/violet/vert), bordures critical/strategic, flèches dépendances."""
    from datetime import datetime as dt, timedelta
    import matplotlib.dates as mdates

    RAG_HEX = {"green": "#10B981", "orange": "#F59E0B", "red": "#EF4444"}
    FAMILY_FILL = {
        "epic_lifecycle": "#EAB308",  # or
        "epic_milestone": "#8B5CF6",  # violet
        "transversal":    "#10B981",  # vert
    }
    IMPACT_COLORS = {"critical": "#EF4444", "high": "#F97316", "medium": "#F59E0B", "low": "#10B981"}

    def _parse(d):
        if not d:
            return None
        try:
            return dt.strptime(str(d)[:10], "%Y-%m-%d")
        except Exception:
            return None

    rows = []
    project_ids_in_slide = set()
    for p in projects:
        start = _parse(p.get("start_date"))
        end   = _parse(p.get("end_date_forecast") or p.get("end_date_baseline"))
        if not start:
            continue
        if not end or end <= start:
            end = start + timedelta(days=90)
        rows.append({
            "name": (p.get("name") or "?")[:35],
            "start": start, "end": end,
            "rag": p.get("status_rag", "green"),
            "pid": p.get("project_id"),
        })
        project_ids_in_slide.add(p.get("project_id"))

    slide = _blank_slide(prs)
    _rect(slide, Emu(0), Emu(0), SW, SH, fill=WHITE)
    _header(slide, "Roadmap Portefeuille",
            subtitle=f"{instance_name} · {instance_date}", brand=brand)
    _footer(slide, f"{(brand or {}).get('company_name', 'MARCEL')} · Confidentiel")

    if not rows:
        tb = _tb(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = "Aucun projet avec des dates planifiées définies."
        r.font.size = Pt(12)
        r.font.name = FONT
        r.font.color.rgb = MID
        return

    t_min = min(r["start"] for r in rows) - timedelta(days=20)
    t_max = max(r["end"]   for r in rows) + timedelta(days=20)
    n     = len(rows)

    fig_h = max(3.0, min(6.5, 1.0 + n * 0.5))
    fig, ax = plt.subplots(figsize=(12.5, fig_h))
    ax.set_facecolor("#FAFAFA")
    fig.patch.set_facecolor("#FAFAFA")

    ax.set_xlim(mdates.date2num(t_min), mdates.date2num(t_max))
    ax.set_ylim(-0.6, n - 0.4)
    ax.invert_yaxis()

    # Grid
    ax.xaxis.set_major_locator(mdates.MonthLocator())
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%b %y"))
    ax.xaxis.set_tick_params(labelsize=7, rotation=25)
    ax.grid(axis="x", linestyle="--", linewidth=0.4, alpha=0.5, color="#CBD5E1")
    ax.tick_params(axis="y", length=0, labelsize=7.5)
    for spine in ["top", "right", "left"]:
        ax.spines[spine].set_visible(False)

    # Today line
    today = dt.now()
    if t_min <= today <= t_max:
        ax.axvline(mdates.date2num(today), color="#0052CC",
                   linewidth=1.5, linestyle="--", alpha=0.6, zorder=5)
        ax.text(mdates.date2num(today), -0.4, "Auj.",
                fontsize=6, color="#0052CC", ha="center", va="bottom", zorder=6)

    # Milestones by project
    ms_by_proj = {}
    for ms in all_milestones:
        pid = ms.get("project_id", "")
        ms_by_proj.setdefault(pid, []).append(ms)

    # Build index: project_id → row index (for dependency arrows)
    pid_to_row = {r["pid"]: i for i, r in enumerate(rows)}

    # Bars & milestones
    for i, r in enumerate(rows):
        color   = RAG_HEX.get(r["rag"], "#10B981")
        x_start = mdates.date2num(r["start"])
        x_end   = mdates.date2num(r["end"])
        bar = mpatches.FancyBboxPatch(
            (x_start, i - 0.28), x_end - x_start, 0.56,
            boxstyle="round,pad=0.3",
            facecolor=color, edgecolor="white", linewidth=0.6, zorder=4,
        )
        ax.add_patch(bar)
        # Label inside bar
        bar_days = (r["end"] - r["start"]).days
        if bar_days > 25:
            cx = (x_start + x_end) / 2
            ax.text(cx, i, r["name"], ha="center", va="center",
                    fontsize=6.5, color="white", fontweight="bold", zorder=5,
                    clip_on=True)

        # Milestone diamonds — colored by family
        for ms in ms_by_proj.get(r["pid"], []):
            d = _parse(ms.get("date_forecast") or ms.get("date_baseline"))
            if not d or not (t_min <= d <= t_max):
                continue
            family = ms.get("family")
            ms_fill = FAMILY_FILL.get(family, "#0B2545")
            attribute = ms.get("attribute")
            edge_color = "#EF4444" if attribute == "critical" else "#3B82F6" if attribute == "strategic" else "white"
            edge_w = 1.5 if attribute else 0.5
            ax.plot(mdates.date2num(d), i, "D",
                    color=ms_fill,
                    markersize=7,
                    markeredgewidth=edge_w,
                    markeredgecolor=edge_color,
                    zorder=7)
            # Bloquant indicator
            if ms.get("is_blocking"):
                ax.text(mdates.date2num(d), i - 0.35, "⚑",
                        fontsize=5, color="#EF4444", ha="center", zorder=8)

    # Dependency arrows between projects in this slide
    if dependencies:
        for dep in dependencies:
            sp = dep.get("source_project_id")
            tp = dep.get("target_project_id")
            if sp not in project_ids_in_slide or tp not in project_ids_in_slide:
                continue
            si = pid_to_row.get(sp)
            ti = pid_to_row.get(tp)
            if si is None or ti is None:
                continue
            src_row = rows[si]
            tgt_row = rows[ti]
            sx = mdates.date2num(src_row["end"])
            tx = mdates.date2num(tgt_row["start"])
            color = IMPACT_COLORS.get(dep.get("impact", "medium"), "#8B5CF6")
            ax.annotate("",
                xy=(tx, ti), xytext=(sx, si),
                arrowprops=dict(
                    arrowstyle="-|>",
                    color=color,
                    lw=1.0,
                    linestyle="dashed",
                    connectionstyle="arc3,rad=0.2",
                ),
                zorder=3)

    ax.set_yticks(range(n))
    ax.set_yticklabels([r["name"] for r in rows], fontsize=7)

    # Legend
    legend_patches = [
        mpatches.Patch(facecolor="#EAB308", label="Epic Lifecycle"),
        mpatches.Patch(facecolor="#8B5CF6", label="Epic Milestone"),
        mpatches.Patch(facecolor="#10B981", label="Transversal / Réglementaire"),
        mpatches.Patch(facecolor="#EAB308", edgecolor="#EF4444", linewidth=2, label="Critical"),
        mpatches.Patch(facecolor="#8B5CF6", edgecolor="#3B82F6", linewidth=2, label="Strategic"),
    ]
    ax.legend(handles=legend_patches, loc="lower right", fontsize=5.5, framealpha=0.8,
              ncol=len(legend_patches), borderpad=0.5)

    ax.set_title("Roadmap Portefeuille", fontsize=10, fontweight="bold",
                 color="#0F172A", pad=6, loc="left")

    plt.tight_layout(pad=0.6)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)

    img_w  = Inches(12.2)
    img_h  = Inches(5.6)
    img_l  = (SW - img_w) / 2
    img_t  = Inches(1.3)
    slide.shapes.add_picture(buf, img_l, img_t, img_w, img_h)


# ---- Header bar (full-width, couleur primaire) ----

def _header(slide, title, subtitle=None, height_in=1.15, brand=None):
    primary = (brand or {}).get("primary", NAVY)
    font = (brand or _CURRENT_BRAND or {}).get("font", FONT)
    h = Inches(height_in)
    # Pleine largeur exacte
    _rect(slide, Emu(0), Emu(0), SW, h, fill=primary)
    tb = _tb(slide, Inches(0.4), Inches(0.12), SW - Inches(0.8), h - Inches(0.25))
    _clear(tb.text_frame)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.name = font
    r.font.color.rgb = WHITE
    if subtitle:
        _run(tb.text_frame, subtitle, size=9, color=RGBColor(0xAA, 0xCC, 0xFF))


# ---- Section label ----

def _section_label(slide, left, top, width, height_in, label):
    h = Inches(height_in)
    _rect(slide, left, top, width, h, fill=LIGHT_BLUE, no_line=False, line_color=BORDER, line_pt=0.3)
    tb = _tb(slide, left + Inches(0.12), top + Inches(0.04), width - Inches(0.2), h)
    _clear(tb.text_frame)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = label.upper()
    r.font.size = Pt(7)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = BLUE


# ---- Slide footer ----

def _footer(slide, text="MARCEL · Confidentiel"):
    tb = _tb(slide, Inches(0.4), SH - Inches(0.3), SW - Inches(1.8), Inches(0.25))
    _clear(tb.text_frame)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(6.5)
    r.font.name = _CURRENT_BRAND.get("font", FONT) if _CURRENT_BRAND else FONT
    r.font.color.rgb = LIGHT
    # Logo tenant en bas à droite (sur toutes les slides)
    logo_b64 = _CURRENT_BRAND.get("logo_base64") if _CURRENT_BRAND else None
    if logo_b64:
        try:
            logo_data = base64.b64decode(logo_b64)
            logo_buf = io.BytesIO(logo_data)
            logo_w = Inches(0.7)
            slide.shapes.add_picture(
                logo_buf, SW - logo_w - Inches(0.15), SH - Inches(0.35), width=logo_w
            )
        except Exception:
            pass


# ---- Table helper (with word wrap) ----

def _styled_table(slide, headers, rows, left, top, width, col_widths_in, row_height_in=0.42,
                  header_bg=NAVY, header_fg=WHITE):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    h_total = Inches(row_height_in) * n_rows
    tf_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, h_total)
    tf = tf_shape.table

    for i, w in enumerate(col_widths_in):
        tf.columns[i].width = Inches(w)

    for row in tf.rows:
        row.height = Inches(row_height_in)

    # Header row
    for col, h in enumerate(headers):
        cell = tf.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.02)
        cell.text_frame.word_wrap = True
        cell.text_frame.paragraphs[0].clear()
        r = cell.text_frame.paragraphs[0].add_run()
        r.text = h
        r.font.size = Pt(8)
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = header_fg
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        row_num = row_idx + 1
        row_bg = BG if row_idx % 2 == 0 else WHITE
        for col, (val, opts) in enumerate(row_data):
            cell = tf.cell(row_num, col)
            cell.fill.solid()
            bg = opts.get("bg", row_bg)
            cell.fill.fore_color.rgb = bg
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.text_frame.word_wrap = True
            cell.text_frame.paragraphs[0].clear()
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = str(val)
            run.font.size = Pt(opts.get("size", 8))
            run.font.bold = opts.get("bold", False)
            run.font.name = FONT
            run.font.color.rgb = opts.get("color", DARK)
            cell.text_frame.paragraphs[0].alignment = opts.get("align", PP_ALIGN.LEFT)

    return tf


# ---- 2-column key/value table helper ----

def _kv_table(slide, items, left, top, width, label_ratio=0.45, row_h=0.36):
    """items = list of (label, value, value_color, bold, row_bg)"""
    n = len(items)
    height = Inches(row_h * n)
    ts = slide.shapes.add_table(n, 2, left, top, width, height)
    tbl = ts.table
    lw = int(width * label_ratio)
    tbl.columns[0].width = lw
    tbl.columns[1].width = width - lw
    for i, (lbl, val, fg, bold, bg) in enumerate(items):
        tbl.rows[i].height = Inches(row_h)
        for col_idx, (text, clr, is_bold) in enumerate([(lbl, MID, False), (val, fg, bold)]):
            cell = tbl.cell(i, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.margin_left = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.text_frame.word_wrap = True
            p = cell.text_frame.paragraphs[0]
            p.clear()
            r = p.add_run()
            r.text = text
            r.font.size = Pt(8.5)
            r.font.bold = is_bold
            r.font.name = FONT
            r.font.color.rgb = clr
    return ts, Inches(row_h * n)


# ====================================================================
# SLIDE 1 — GARDE (fond blanc, titre Navy, bloc projets encadré)
# ====================================================================

def add_slide_garde(prs, instance_name, instance_date, projects, brand=None):
    br = brand or {}
    primary    = br.get("primary", NAVY)
    co_name    = br.get("company_name", "MARCEL")
    logo_b64   = br.get("logo_base64")

    slide = _blank_slide(prs)

    # Thin top accent bar (couleur primaire)
    _rect(slide, Emu(0), Emu(0), SW, Inches(0.1), fill=primary)

    # App label
    app_tb = _tb(slide, Inches(0.5), Inches(0.22), Inches(8), Inches(0.32))
    _clear(app_tb.text_frame)
    r = app_tb.text_frame.paragraphs[0].add_run()
    r.text = f"{co_name.upper()}  ·  EXPORT COPIL"
    r.font.size = Pt(8)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = primary

    # Logo tenant (haut à droite, si configuré)
    if logo_b64:
        try:
            logo_data = base64.b64decode(logo_b64)
            logo_buf  = io.BytesIO(logo_data)
            logo_w    = Inches(1.5)
            slide.shapes.add_picture(
                logo_buf, SW - logo_w - Inches(0.3), Inches(0.1), width=logo_w
            )
        except Exception:
            pass

    # Instance name — large, couleur primaire
    ttb = _tb(slide, Inches(0.5), Inches(0.65), SW - Inches(1.0), Inches(1.5))
    ttb.text_frame.word_wrap = True
    _clear(ttb.text_frame)
    p2 = ttb.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = instance_name
    r2.font.size = Pt(36)
    r2.font.bold = True
    r2.font.name = FONT
    r2.font.color.rgb = primary

    # Date — dark text
    dtb = _tb(slide, Inches(0.5), Inches(2.3), Inches(8), Inches(0.4))
    _clear(dtb.text_frame)
    r3 = dtb.text_frame.paragraphs[0].add_run()
    r3.text = fmt_date(instance_date)
    r3.font.size = Pt(14)
    r3.font.name = FONT
    r3.font.color.rgb = DARK

    # Scope count — dark text
    scope_tb = _tb(slide, Inches(0.5), Inches(2.78), SW - Inches(1.0), Inches(0.35))
    _clear(scope_tb.text_frame)
    r4 = scope_tb.text_frame.paragraphs[0].add_run()
    r4.text = f"{len(projects)} projet{'s' if len(projects) > 1 else ''} en périmètre"
    r4.font.size = Pt(11)
    r4.font.name = FONT
    r4.font.color.rgb = MID

    # Header bar pour liste projets (couleur primaire)
    list_start = Inches(3.25)
    _rect(slide, Inches(0.4), list_start, SW - Inches(0.8), Inches(0.3), fill=primary)
    hdr_tb = _tb(slide, Inches(0.55), list_start + Inches(0.05), SW - Inches(1.1), Inches(0.25))
    _clear(hdr_tb.text_frame)
    rh = hdr_tb.text_frame.paragraphs[0].add_run()
    rh.text = "PROJETS EN PÉRIMÈTRE"
    rh.font.size = Pt(7.5)
    rh.font.bold = True
    rh.font.name = FONT
    rh.font.color.rgb = WHITE

    # Project rows — encadrés
    row_h = Inches(0.42)
    for idx, p_data in enumerate(projects[:8]):
        row_y = list_start + Inches(0.3) + idx * row_h
        rag = p_data.get("status_rag", "green")
        _rect(slide, Inches(0.4), row_y, SW - Inches(0.8), row_h - Inches(0.03),
              fill=BG, no_line=False, line_color=BORDER, line_pt=0.3)
        _rect(slide, Inches(0.4), row_y, Inches(0.07), row_h - Inches(0.03), fill=rag_color(rag))
        ntb = _tb(slide, Inches(0.58), row_y + Inches(0.07), Inches(10.0), row_h)
        _clear(ntb.text_frame)
        rn = ntb.text_frame.paragraphs[0].add_run()
        rn.text = trunc(p_data.get("name", "?"), 65)
        rn.font.size = Pt(9)
        rn.font.bold = True
        rn.font.name = FONT
        rn.font.color.rgb = DARK
        btb = _tb(slide, SW - Inches(2.4), row_y + Inches(0.07), Inches(2.0), row_h)
        _clear(btb.text_frame)
        pp2 = btb.text_frame.paragraphs[0]
        pp2.alignment = PP_ALIGN.RIGHT
        rb = pp2.add_run()
        rb.text = keur(p_data.get("budget_total"))
        rb.font.size = Pt(8)
        rb.font.name = FONT
        rb.font.color.rgb = MID

    if len(projects) > 8:
        etb = _tb(slide, Inches(0.5),
                  list_start + Inches(0.3) + 8 * row_h + Inches(0.06),
                  Inches(8), Inches(0.28))
        _clear(etb.text_frame)
        re = etb.text_frame.paragraphs[0].add_run()
        re.text = f"+ {len(projects) - 8} autre(s) projet(s)…"
        re.font.size = Pt(8)
        re.font.name = FONT
        re.font.color.rgb = LIGHT

    _footer(slide, f"{co_name} · Confidentiel")


# ====================================================================
# SLIDE 2 — SOMMAIRE (tableau pleine largeur)
# ====================================================================

def add_slide_sommaire(prs, projects, instance_name="", instance_date="", brand=None):
    slide = _blank_slide(prs)
    _header(slide, "Synthèse du Portefeuille",
            f"{len(projects)} projet{'s' if len(projects) > 1 else ''} sélectionné{'s' if len(projects) > 1 else ''}",
            brand=brand)

    headers = ["RAG", "Projet", "Responsable", "Budget (K€)", "EAC (K€)", "Écart", "Statut", "Fin forecast"]
    # Colonnes élargies pour couvrir toute la largeur (~12.93")
    col_w = [0.6, 4.1, 1.8, 1.3, 1.3, 1.1, 1.4, 1.33]
    rows = []

    for p in projects:
        total = p.get("budget_total", 0) or 0
        eac = p.get("eac") or total
        ecart = eac - total
        ecart_pct = (ecart / total * 100) if total else 0
        ecart_str = f"{'+'if ecart>0 else ''}{int(ecart/1000):,}".replace(",", "\u202f")
        ecart_bg = LIGHT_RED if ecart > 0 else LIGHT_GREEN if ecart < 0 else BG
        rag = p.get("status_rag", "green")

        rows.append([
            (rag_label(rag), {"bg": {"green": LIGHT_GREEN, "orange": LIGHT_AMBER, "red": LIGHT_RED}.get(rag, BG),
                              "color": rag_color(rag), "bold": True, "align": PP_ALIGN.CENTER, "size": 8}),
            (trunc(p.get("name", "?"), 55), {}),
            (trunc(p.get("owner_name", p.get("metadata", {}).get("sponsor", "—")), 22), {"color": MID}),
            (f"{int(total/1000):,}".replace(",", "\u202f"), {"align": PP_ALIGN.RIGHT}),
            (f"{int(eac/1000):,}".replace(",", "\u202f"), {"align": PP_ALIGN.RIGHT, "bold": True}),
            (ecart_str + " K€", {"align": PP_ALIGN.RIGHT, "bg": ecart_bg,
                                  "color": RED if ecart > 0 else GREEN_C if ecart < 0 else DARK}),
            (status_label(p.get("status")), {"color": MID, "size": 7}),
            (fmt_date(p.get("end_date_forecast")), {"align": PP_ALIGN.CENTER, "color": MID}),
        ])

    table_left = Inches(0.2)
    table_width = sum(Inches(w) for w in col_w)
    _styled_table(slide, headers, rows, table_left, Inches(1.25), table_width, col_w, row_height_in=0.42)

    co = (brand or {}).get("company_name", "MARCEL")
    footer_text = f"{instance_name}  ·  {fmt_date(instance_date)}  ·  {co} Confidentiel" if instance_name else f"{co} · Confidentiel"
    _footer(slide, footer_text)


# ====================================================================
# SLIDE 3 — HEATMAP P×I (axe X corrigé)
# ====================================================================

def generate_heatmap_img(risks):
    fig, ax = plt.subplots(figsize=(5.5, 5.2))
    ax.set_facecolor("white")
    fig.patch.set_facecolor("white")

    for p in range(1, 6):
        for i in range(1, 6):
            crit = p * i
            fc = "#FEF2F2" if crit >= 16 else "#FFF7ED" if crit >= 7 else "#F0FDF4"
            ec = "#FECACA" if crit >= 16 else "#FED7AA" if crit >= 7 else "#BBF7D0"
            rect = mpatches.FancyBboxPatch(
                (i - 0.48, p - 0.48), 0.96, 0.96,
                boxstyle="square,pad=0", facecolor=fc, edgecolor=ec, linewidth=0.7, zorder=1,
            )
            ax.add_patch(rect)
            ax.text(i, p, str(crit), ha="center", va="center",
                    fontsize=8, color="#CBD5E1", zorder=2)

    risk_pos: dict = {}
    for r in risks:
        k = (int(r.get("impact", 1)), int(r.get("probability", 1)))
        risk_pos[k] = risk_pos.get(k, 0) + 1

    for (imp, prob), count in risk_pos.items():
        crit = prob * imp
        color = "#EF4444" if crit >= 16 else "#F59E0B" if crit >= 7 else "#10B981"
        ax.scatter(imp, prob, s=min(180 + count * 40, 450), c=color,
                   zorder=5, edgecolors="white", linewidths=1.5, alpha=0.9)
        if count > 1:
            ax.text(imp, prob, str(count), ha="center", va="center",
                    fontsize=8, color="white", fontweight="bold", zorder=6)

    ax.set_xlim(0.48, 5.52)
    ax.set_ylim(0.48, 5.52)
    ax.set_xticks(range(1, 6))
    ax.set_yticks(range(1, 6))
    ax.set_xticklabels([str(i) for i in range(1, 6)], fontsize=9, color="#475569")
    ax.set_yticklabels([str(i) for i in range(1, 6)], fontsize=9, color="#475569")
    ax.set_xlabel("Impact  →", fontsize=10, color="#334155", labelpad=8)
    ax.set_ylabel("Probabilité  →", fontsize=10, color="#334155", labelpad=8)
    ax.tick_params(length=0)
    for spine in ax.spines.values():
        spine.set_edgecolor("#E2E8F0")
        spine.set_linewidth(0.5)
    ax.set_aspect("equal")
    # FIX: réserver espace en bas pour l'axe X (évite la coupure)
    plt.subplots_adjust(left=0.14, right=0.96, bottom=0.18, top=0.96)

    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=150, facecolor="white", edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


def add_slide_heatmap(prs, risks, proj_names, instance_name="", instance_date="", brand=None):
    slide = _blank_slide(prs)
    crit_n = sum(1 for r in risks if r.get("criticality", 0) >= 16)
    mod_n = sum(1 for r in risks if 7 <= r.get("criticality", 0) < 16)
    low_n = sum(1 for r in risks if r.get("criticality", 0) < 7)
    _header(slide, "Cartographie des Risques P × I",
            f"{len(risks)} risques au total  ·  Élevés : {crit_n}  ·  Modérés : {mod_n}  ·  Faibles : {low_n}",
            brand=brand)

    img_buf = generate_heatmap_img(risks)
    slide.shapes.add_picture(img_buf, Inches(0.3), Inches(1.25), width=Inches(6.7))

    rx = Inches(7.3)
    ry = Inches(1.4)
    for label, count, bg, fg in [
        ("Risques ÉLEVÉS (P×I ≥ 16)", crit_n, LIGHT_RED, RED),
        ("Risques MODÉRÉS (P×I 7–15)", mod_n, LIGHT_AMBER, ORANGE_C),
        ("Risques FAIBLES (P×I ≤ 6)", low_n, LIGHT_GREEN, GREEN_C),
    ]:
        _rect(slide, rx, ry, Inches(5.7), Inches(0.75), fill=bg, no_line=False, line_color=BORDER, line_pt=0.3)
        tb_l = _tb(slide, rx + Inches(0.15), ry + Inches(0.08), Inches(4.0), Inches(0.6))
        _clear(tb_l.text_frame)
        _run(tb_l.text_frame, label, size=8, color=DARK)
        tb_n = _tb(slide, rx + Inches(4.2), ry + Inches(0.1), Inches(1.3), Inches(0.5))
        _clear(tb_n.text_frame)
        p = tb_n.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = str(count)
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = fg
        ry += Inches(0.85)

    top_risks = sorted(risks, key=lambda x: -x.get("criticality", 0))[:5]
    if top_risks:
        ry += Inches(0.15)
        lbl_tb = _tb(slide, rx, ry, Inches(5.7), Inches(0.26))
        _clear(lbl_tb.text_frame)
        _run(lbl_tb.text_frame, "TOP 5 RISQUES CRITIQUES", size=7, bold=True, color=MID)
        ry += Inches(0.28)
        for r_item in top_risks:
            crit = r_item.get("criticality", 0)
            pname = trunc(proj_names.get(r_item.get("project_id", ""), "?"), 20)
            _rect(slide, rx, ry, Inches(5.7), Inches(0.44), fill=BG, no_line=False, line_color=BORDER, line_pt=0.3)
            _rect(slide, rx + Inches(0.08), ry + Inches(0.12), Inches(0.2), Inches(0.2), fill=crit_color(crit))
            tb_r = _tb(slide, rx + Inches(0.38), ry + Inches(0.07), Inches(5.2), Inches(0.35))
            _clear(tb_r.text_frame)
            p = tb_r.text_frame.paragraphs[0]
            r_run = p.add_run()
            r_run.text = f"[{crit}] {trunc(r_item.get('title', '?'), 45)}"
            r_run.font.size = Pt(7.5)
            r_run.font.name = FONT
            r_run.font.color.rgb = DARK
            p2 = tb_r.text_frame.add_paragraph()
            r2 = p2.add_run()
            r2.text = pname
            r2.font.size = Pt(6.5)
            r2.font.name = FONT
            r2.font.color.rgb = MID
            ry += Inches(0.47)

    footer_text = f"{instance_name}  ·  {fmt_date(instance_date)}  ·  {(brand or {}).get('company_name', 'MARCEL')} Confidentiel" if instance_name else f"{(brand or {}).get('company_name', 'MARCEL')} · Confidentiel"
    _footer(slide, footer_text)


# ====================================================================
# SLIDE 4 — TOP RISQUES CRITIQUES (colonnes élargies, word wrap)
# ====================================================================

def add_slide_top_risks(prs, risks, proj_names, instance_name="", instance_date="", brand=None):
    slide = _blank_slide(prs)
    top = sorted(risks, key=lambda x: -x.get("criticality", 0))[:10]
    _header(slide, "Top Risques Critiques",
            f"Projets sélectionnés — {len(top)} risques présentés par criticité décroissante",
            brand=brand)

    headers = ["#", "Crit.", "Risque", "Catégorie", "Projet", "Statut", "Échéance", "Propriétaire"]
    # Colonnes élargies : Risque 3.5→4.3, Projet 2.3→3.0
    col_w = [0.35, 0.5, 4.3, 1.1, 3.0, 0.85, 1.0, 1.83]
    rows = []
    for idx, r in enumerate(top):
        crit = r.get("criticality", 0)
        rows.append([
            (str(idx + 1), {"align": PP_ALIGN.CENTER, "color": MID}),
            (str(crit), {"align": PP_ALIGN.CENTER, "bold": True, "color": crit_color(crit), "bg": crit_bg(crit)}),
            (trunc(r.get("title", "?"), 70), {}),
            (r.get("category", "—"), {"color": MID}),
            (trunc(proj_names.get(r.get("project_id", ""), "?"), 42), {"color": BLUE}),
            (r.get("status", "—"), {"color": MID, "size": 7}),
            (fmt_date(r.get("due_date")), {"align": PP_ALIGN.CENTER, "color": MID}),
            (trunc(r.get("owner", "—") or "—", 25), {"color": MID}),
        ])

    left = Inches(0.2)
    width = sum(Inches(w) for w in col_w)
    _styled_table(slide, headers, rows, left, Inches(1.25), width, col_w, row_height_in=0.46)

    footer_text = f"{instance_name}  ·  {fmt_date(instance_date)}  ·  {(brand or {}).get('company_name', 'MARCEL')} Confidentiel" if instance_name else f"{(brand or {}).get('company_name', 'MARCEL')} · Confidentiel"
    _footer(slide, footer_text)


# ====================================================================
# SLIDE 5 — DÉCISIONS CLÉS (colonnes élargies, footer)
# ====================================================================

def add_slide_decisions(prs, decisions, governance_id=None, instance_name="", instance_date="", brand=None):
    slide = _blank_slide(prs)
    subtitle = "Instance de gouvernance sélectionnée" if governance_id else "10 dernières décisions du périmètre"
    _header(slide, "Décisions Clés",
            f"{len(decisions)} décision{'s' if len(decisions) != 1 else ''}  ·  {subtitle}",
            brand=brand)

    headers = ["Date", "Décision", "Catégorie", "Statut", "Projet", "Responsable"]
    # Décision 4.0→5.0, Projet 3.0→3.0 (keep), Responsable ajusté
    col_w = [0.95, 5.0, 1.2, 1.0, 3.0, 1.78]
    rows = []
    for d in decisions[:12]:
        status = d.get("status", "—")
        rows.append([
            (fmt_date(d.get("decision_date")), {"align": PP_ALIGN.CENTER, "color": MID}),
            (trunc(d.get("title", "?"), 80), {}),
            (d.get("category", "—"), {"color": MID}),
            (decision_status_label(status), {
                "align": PP_ALIGN.CENTER, "size": 7,
                "bg": DECISION_STATUS_COLORS.get(status, BG), "bold": True,
            }),
            (trunc(d.get("project_name", "?"), 42), {"color": BLUE}),
            (trunc(d.get("owner", "—") or "—", 25), {"color": MID}),
        ])

    left = Inches(0.2)
    width = sum(Inches(w) for w in col_w)
    _styled_table(slide, headers, rows, left, Inches(1.25), width, col_w, row_height_in=0.46)

    footer_text = f"{instance_name}  ·  {fmt_date(instance_date)}  ·  {(brand or {}).get('company_name', 'MARCEL')} Confidentiel" if instance_name else f"{(brand or {}).get('company_name', 'MARCEL')} · Confidentiel"
    _footer(slide, footer_text)


# ====================================================================
# SLIDE 6+N — FICHE PROJET (tables KV pour budget/planning, sous-blocs risques/décisions)
# ====================================================================

def add_slide_fiche(prs, project, milestones, risks, decisions, instance_name="", instance_date="", brand=None):
    slide = _blank_slide(prs)
    primary = (brand or {}).get("primary", NAVY)
    co      = (brand or {}).get("company_name", "MARCEL")

    name = project.get("name", "?")
    code = project.get("source_id", "")
    rag = project.get("status_rag", "green")
    status = project.get("status", "")
    methodology = project.get("methodology", "—")
    owner = project.get("owner_name") or project.get("metadata", {}).get("sponsor", "—") or "—"

    # ---- Header ----
    HEADER_H = Inches(1.2)
    _rect(slide, Emu(0), Emu(0), SW, HEADER_H, fill=primary)
    _rect(slide, Emu(0), Emu(0), Inches(0.18), HEADER_H, fill=rag_color(rag))

    ttb = _tb(slide, Inches(0.32), Inches(0.1), SW - Inches(4.0), Inches(0.65))
    _clear(ttb.text_frame)
    r = ttb.text_frame.paragraphs[0].add_run()
    r.text = trunc(name, 72)
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = WHITE

    ctb = _tb(slide, Inches(0.32), Inches(0.78), Inches(4), Inches(0.32))
    _clear(ctb.text_frame)
    r2 = ctb.text_frame.paragraphs[0].add_run()
    r2.text = code if code else ""
    r2.font.size = Pt(8)
    r2.font.name = FONT
    r2.font.color.rgb = LIGHT

    meta_tb = _tb(slide, Inches(4.0), Inches(0.78), SW - Inches(4.4), Inches(0.32))
    _clear(meta_tb.text_frame)
    p3 = meta_tb.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.RIGHT
    r3 = p3.add_run()
    r3.text = f"{owner}  ·  {methodology.upper()}  ·  {status_label(status)}  ·  RAG {rag_label(rag).upper()}"
    r3.font.size = Pt(8)
    r3.font.name = FONT
    r3.font.color.rgb = rag_color(rag)

    # ---- Layout grid ----
    BODY_Y = HEADER_H + Inches(0.08)
    COL_W = Inches(6.35)
    L_X = Inches(0.2)
    R_X = SW - COL_W - Inches(0.2)
    TOP_H = Inches(2.75)
    BTM_Y = BODY_Y + TOP_H + Inches(0.1)
    BTM_H = SH - BTM_Y - Inches(0.35)

    # ---- BUDGET — table KV 2 colonnes (top-left) ----
    capex_p = project.get("capex_planned", 0) or 0
    capex_c = project.get("capex_consumed", 0) or 0
    opex_p = project.get("opex_planned", 0) or 0
    opex_c = project.get("opex_consumed", 0) or 0
    total = project.get("budget_total", 0) or 0
    consumed = project.get("budget_consumed", 0) or (capex_c + opex_c)
    eac = project.get("eac") or total
    ecart = eac - total
    ecart_pct = (ecart / total * 100) if total else 0
    cons_pct = min(int(consumed / total * 100), 100) if total else 0
    eac_col = RED if ecart > 0 else GREEN_C if ecart < 0 else DARK
    ecart_str = (f"{'+'if ecart>0 else ''}{int(ecart/1000):,}".replace(",", "\u202f") +
                 f" K€  ({'+' if ecart_pct > 0 else ''}{ecart_pct:.1f}%)" if total else "—")

    _section_label(slide, L_X, BODY_Y, COL_W, 0.26, "Budget")
    bud_y = BODY_Y + Inches(0.26)

    budget_items = [
        ("CAPEX prévu",    keur(capex_p),  DARK,   False, BG),
        ("CAPEX consommé", keur(capex_c),  DARK,   True,  WHITE),
        ("OPEX prévu",     keur(opex_p),   DARK,   False, BG),
        ("OPEX consommé",  keur(opex_c),   DARK,   True,  WHITE),
        ("Budget total",   keur(total),    DARK,   True,  BG),
        ("EAC (Estimate At Completion)", keur(eac), eac_col, True, LIGHT_BLUE),
        ("Écart EAC / Budget", ecart_str,  eac_col, True, LIGHT_RED if ecart > 0 else LIGHT_GREEN if ecart < 0 else BG),
    ]
    _, kv_h = _kv_table(slide, budget_items, L_X, bud_y, COL_W, label_ratio=0.45, row_h=0.33)
    bud_y += kv_h + Inches(0.06)

    # Progress bar
    _rect(slide, L_X, bud_y, COL_W, Inches(0.13), fill=BORDER)
    if cons_pct > 0:
        bar_w = int(COL_W * cons_pct / 100)
        _rect(slide, L_X, bud_y, bar_w, Inches(0.13),
              fill=RED if cons_pct >= 90 else BLUE)
    bud_y += Inches(0.15)
    pct_tb = _tb(slide, L_X, bud_y, COL_W, Inches(0.2))
    _clear(pct_tb.text_frame)
    r_pct = pct_tb.text_frame.paragraphs[0].add_run()
    r_pct.text = f"Avancement budgétaire : {cons_pct}%  ({keur(consumed)} / {keur(total)})"
    r_pct.font.size = Pt(7)
    r_pct.font.name = FONT
    r_pct.font.color.rgb = MID

    # ---- PLANNING — table KV 2 colonnes (top-right) ----
    start = project.get("start_date", "") or ""
    end_base = project.get("end_date_baseline", "") or ""
    end_fore = project.get("end_date_forecast", "") or ""
    delay_str, delay_color = "—", DARK
    if end_base and end_fore:
        try:
            d_base = datetime.fromisoformat(str(end_base).replace("Z", ""))
            d_fore = datetime.fromisoformat(str(end_fore).replace("Z", ""))
            days = (d_fore - d_base).days
            if days > 0:
                delay_str, delay_color = f"+{days} jours (retard)", RED
            elif days < 0:
                delay_str, delay_color = f"{abs(days)} jours (avance)", GREEN_C
            else:
                delay_str, delay_color = "Dans les délais", GREEN_C
        except Exception:
            pass

    _section_label(slide, R_X, BODY_Y, COL_W, 0.26, "Planning")
    plan_y = BODY_Y + Inches(0.26)

    planning_items = [
        ("Début",                        fmt_date(start),    MID,         False, BG),
        ("Fin baseline (initiale)",       fmt_date(end_base), MID,         False, WHITE),
        ("Fin forecast (actuelle)",       fmt_date(end_fore), RED if delay_color == RED else MID, True, BG),
        ("Délai vs baseline",             delay_str,          delay_color, True,  WHITE),
        ("Avancement JH",
         f"{int(project.get('jh_consumed',0)):,} / {int(project.get('jh_planned',0)):,} JH".replace(",", "\u202f"),
         DARK, False, BG),
    ]
    _kv_table(slide, planning_items, R_X, plan_y, COL_W, label_ratio=0.48, row_h=0.38)

    # ---- JALONS (bottom-left) ----
    _section_label(slide, L_X, BTM_Y, COL_W, 0.26, "Jalons Clés")
    mil_y = BTM_Y + Inches(0.26)
    now_str = datetime.now().isoformat()
    past = [m for m in milestones if (m.get("date_forecast") or "") < now_str][-3:]
    future = [m for m in milestones if (m.get("date_forecast") or "") >= now_str][:3]
    combined = [(False, m) for m in reversed(past)] + [(True, m) for m in future]

    for i, (upcoming, ms) in enumerate(combined[:6]):
        my = mil_y + i * Inches(0.4)
        _rect(slide, L_X, my, COL_W, Inches(0.38),
              fill=LIGHT_BLUE if upcoming else BG,
              no_line=False, line_color=BORDER, line_pt=0.3)
        dtb = _tb(slide, L_X + Inches(0.1), my + Inches(0.07), Inches(1.15), Inches(0.28))
        _clear(dtb.text_frame)
        _run(dtb.text_frame, fmt_date(ms.get("date_forecast")), size=7.5,
             color=BLUE if upcoming else MID, bold=upcoming)
        ntb = _tb(slide, L_X + Inches(1.3), my + Inches(0.07), COL_W - Inches(1.45), Inches(0.28))
        _clear(ntb.text_frame)
        _run(ntb.text_frame, trunc(ms.get("name", "?"), 44), size=7.5, color=DARK)

    if not combined:
        ntb = _tb(slide, L_X + Inches(0.1), mil_y + Inches(0.1), COL_W - Inches(0.2), Inches(0.3))
        _clear(ntb.text_frame)
        _run(ntb.text_frame, "Aucun jalon défini.", size=8, color=LIGHT, italic=True)

    # ---- RISQUES & DÉCISIONS (bottom-right) — 2 sous-blocs séparés ----
    rd_y = BTM_Y

    # Sous-bloc RISQUES
    _section_label(slide, R_X, rd_y, COL_W, 0.26, f"Risques (Top {min(3, len(risks))})")
    rd_y += Inches(0.26)
    top3 = sorted(risks, key=lambda x: -x.get("criticality", 0))[:3]
    for r_item in top3:
        crit = r_item.get("criticality", 0)
        _rect(slide, R_X, rd_y, COL_W, Inches(0.42), fill=crit_bg(crit),
              no_line=False, line_color=BORDER, line_pt=0.3)
        _rect(slide, R_X + Inches(0.1), rd_y + Inches(0.13),
              Inches(0.18), Inches(0.18), fill=crit_color(crit))
        rtb = _tb(slide, R_X + Inches(0.38), rd_y + Inches(0.06),
                  COL_W - Inches(0.55), Inches(0.34))
        _clear(rtb.text_frame)
        _run(rtb.text_frame, f"[{crit}] {trunc(r_item.get('title', '?'), 52)}", size=7.5, color=DARK)
        rd_y += Inches(0.44)
    if not top3:
        etb = _tb(slide, R_X + Inches(0.1), rd_y + Inches(0.05), COL_W, Inches(0.3))
        _clear(etb.text_frame)
        _run(etb.text_frame, "Aucun risque enregistré.", size=8, color=LIGHT, italic=True)
        rd_y += Inches(0.35)

    rd_y += Inches(0.1)

    # Sous-bloc DÉCISIONS
    _section_label(slide, R_X, rd_y, COL_W, 0.26, f"Décisions ({min(3, len(decisions))} récentes)")
    rd_y += Inches(0.26)
    last3 = decisions[:3]
    for d in last3:
        dstatus = d.get("status", "—")
        _rect(slide, R_X, rd_y, COL_W, Inches(0.42),
              fill=DECISION_STATUS_COLORS.get(dstatus, BG),
              no_line=False, line_color=BORDER, line_pt=0.3)
        dtb2 = _tb(slide, R_X + Inches(0.1), rd_y + Inches(0.06),
                   COL_W - Inches(1.7), Inches(0.34))
        _clear(dtb2.text_frame)
        _run(dtb2.text_frame, trunc(d.get("title", "?"), 46), size=7.5, color=DARK)
        date_tb = _tb(slide, R_X + COL_W - Inches(1.6), rd_y + Inches(0.06),
                      Inches(1.5), Inches(0.34))
        _clear(date_tb.text_frame)
        pp_d = date_tb.text_frame.paragraphs[0]
        pp_d.alignment = PP_ALIGN.RIGHT
        rr = pp_d.add_run()
        rr.text = f"{fmt_date(d.get('decision_date'))}  ·  {decision_status_label(dstatus)[:6]}"
        rr.font.size = Pt(6.5)
        rr.font.name = FONT
        rr.font.color.rgb = MID
        rd_y += Inches(0.44)
    if not last3:
        etb2 = _tb(slide, R_X + Inches(0.1), rd_y + Inches(0.05), COL_W, Inches(0.3))
        _clear(etb2.text_frame)
        _run(etb2.text_frame, "Aucune décision enregistrée.", size=8, color=LIGHT, italic=True)

    # Footer neutre (suppression du footer orphelin)
    footer_text = f"{instance_name}  ·  {fmt_date(instance_date)}  ·  {co} Confidentiel" if instance_name else f"{co} · Confidentiel"
    _footer(slide, footer_text)


# ====================================================================
# SLIDE — CONSOMMATION PAR ÉQUIPE (S1-08)
# ====================================================================

def add_slide_team_consumption(prs, project, team_rows, instance_name, instance_date, brand=None):
    """Slide additionnel par projet : tableau consommation par équipe + KPI RAF."""
    slide = _blank_slide(prs)

    # Background
    _rect(slide, Emu(0), Emu(0), SW, SH, fill=WHITE)

    proj_name = project.get("name", "?")
    _header(slide, f"{proj_name.upper()}  ·  CONSOMMATION PAR ÉQUIPE",
            subtitle="Répartition JH & coûts · RAF valorisé · Atterrissage", height_in=1.2, brand=brand)

    # ---- Compute totals ----
    total_planned_md   = sum(r.get("planned_md", 0) for r in team_rows)
    total_consumed_md  = sum(r.get("consumed_md", 0) for r in team_rows)
    total_raf_md       = sum(r.get("raf_md", 0) for r in team_rows)
    total_planned_eur  = sum(r.get("planned_cost_eur", 0) for r in team_rows)
    total_consumed_eur = sum(r.get("consumed_cost_eur", 0) for r in team_rows)
    total_raf_eur      = sum(r.get("raf_cost_eur", 0) for r in team_rows)
    atterrissage       = total_consumed_eur + total_raf_eur

    # ---- KPI cards row ----
    KPI_Y   = Inches(1.35)
    KPI_H   = Inches(1.25)
    KPI_W   = Inches(3.6)
    KPI_PAD = Inches(0.4)
    kpis = [
        ("Consommé total", eur(total_consumed_eur), f"{jh_fmt(total_consumed_md)}", BLUE),
        ("RAF valorisé",   eur(total_raf_eur),      f"{jh_fmt(total_raf_md)}",      ORANGE_C),
        ("Atterrissage",   eur(atterrissage),        f"prev. {eur(total_planned_eur)}", NAVY),
    ]
    for i, (label, value, sub, accent) in enumerate(kpis):
        kx = KPI_PAD + (KPI_W + Inches(0.3)) * i
        _rect(slide, kx, KPI_Y, KPI_W, KPI_H, fill=RGBColor(0xF8, 0xF9, 0xFA),
              no_line=False, line_color=BORDER, line_pt=0.5)
        # Accent bar (top left)
        _rect(slide, kx, KPI_Y, Inches(0.12), KPI_H, fill=accent)

        lbl_tb = _tb(slide, kx + Inches(0.25), KPI_Y + Inches(0.12), KPI_W - Inches(0.35), Inches(0.3))
        _clear(lbl_tb.text_frame)
        _run(lbl_tb.text_frame, label.upper(), size=7, bold=True, color=MID)

        val_tb = _tb(slide, kx + Inches(0.25), KPI_Y + Inches(0.4), KPI_W - Inches(0.35), Inches(0.5))
        _clear(val_tb.text_frame)
        _run(val_tb.text_frame, value, size=17, bold=True, color=accent)

        sub_tb = _tb(slide, kx + Inches(0.25), KPI_Y + Inches(0.88), KPI_W - Inches(0.35), Inches(0.3))
        _clear(sub_tb.text_frame)
        _run(sub_tb.text_frame, sub, size=8, color=MID)

    # ---- Table consommation par équipe ----
    TABLE_Y = KPI_Y + KPI_H + Inches(0.3)
    TABLE_X = Inches(0.4)
    TABLE_W = SW - Inches(0.8)

    headers = ["Équipe", "JH prévus", "JH consommés", "Coût prévu", "Coût consommé", "RAF JH", "RAF €"]
    # col widths (must sum ≈ TABLE_W / Inches)
    col_w = [2.8, 1.15, 1.35, 1.65, 1.65, 1.1, 1.65]

    rows_data = []
    for row in team_rows:
        raf_md   = row.get("raf_md", 0)
        raf_eur  = row.get("raf_cost_eur", 0)
        rows_data.append([
            (row.get("team_name", "?"),                    {"bold": True, "color": NAVY}),
            (jh_fmt(row.get("planned_md", 0)),             {"align": PP_ALIGN.RIGHT, "color": DARK}),
            (jh_fmt(row.get("consumed_md", 0)),            {"align": PP_ALIGN.RIGHT, "color": DARK, "bold": True}),
            (eur(row.get("planned_cost_eur", 0)),          {"align": PP_ALIGN.RIGHT, "color": DARK}),
            (eur(row.get("consumed_cost_eur", 0)),         {"align": PP_ALIGN.RIGHT, "color": DARK, "bold": True}),
            (jh_fmt(raf_md),                               {"align": PP_ALIGN.RIGHT, "color": ORANGE_C if raf_md > 0 else MID}),
            (eur(raf_eur),                                 {"align": PP_ALIGN.RIGHT, "color": ORANGE_C if raf_eur > 0 else MID, "bold": raf_eur > 0}),
        ])

    # Total row
    rows_data.append([
        ("TOTAL",                    {"bold": True, "color": WHITE, "bg": NAVY}),
        (jh_fmt(total_planned_md),   {"align": PP_ALIGN.RIGHT, "color": WHITE, "bg": NAVY}),
        (jh_fmt(total_consumed_md),  {"align": PP_ALIGN.RIGHT, "color": WHITE, "bg": NAVY, "bold": True}),
        (eur(total_planned_eur),     {"align": PP_ALIGN.RIGHT, "color": WHITE, "bg": NAVY}),
        (eur(total_consumed_eur),    {"align": PP_ALIGN.RIGHT, "color": WHITE, "bg": NAVY, "bold": True}),
        (jh_fmt(total_raf_md),       {"align": PP_ALIGN.RIGHT, "color": ORANGE_C if total_raf_md > 0 else WHITE, "bg": NAVY, "bold": True}),
        (eur(total_raf_eur),         {"align": PP_ALIGN.RIGHT, "color": ORANGE_C if total_raf_eur > 0 else WHITE, "bg": NAVY, "bold": True}),
    ])

    if team_rows:
        _styled_table(slide, headers, rows_data, TABLE_X, TABLE_Y, TABLE_W, col_w, row_height_in=0.45)
    else:
        empty_tb = _tb(slide, TABLE_X, TABLE_Y + Inches(0.3), TABLE_W, Inches(0.5))
        _clear(empty_tb.text_frame)
        _run(empty_tb.text_frame, "Aucune allocation de travail enregistrée pour ce projet.",
             size=9, color=LIGHT, italic=True)

    footer_text = f"{instance_name}  ·  {fmt_date(instance_date)}  ·  {(brand or {}).get('company_name', 'MARCEL')} Confidentiel"
    _footer(slide, footer_text)


# ====================================================================
# SLIDE CLÔTURE — Points d'attention CIO
# ====================================================================

def add_slide_cloture(prs, projects, all_risks, instance_name="", instance_date="", brand=None):
    """Slide finale : points d'attention CIO — risques critiques, projets rouges, EAC dépassés."""
    slide = _blank_slide(prs)
    primary = (brand or {}).get("primary", NAVY)
    co = (brand or {}).get("company_name", "MARCEL")

    # ---- Header ----
    HEADER_H = Inches(1.1)
    _rect(slide, Emu(0), Emu(0), SW, HEADER_H, fill=primary)
    _rect(slide, Emu(0), Emu(0), Inches(0.18), HEADER_H, fill=RED)

    ttb = _tb(slide, Inches(0.32), Inches(0.1), SW - Inches(3.2), Inches(0.55))
    _clear(ttb.text_frame)
    r = ttb.text_frame.paragraphs[0].add_run()
    r.text = "Points d'attention CIO"
    r.font.size = Pt(18); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = WHITE

    sub_tb = _tb(slide, Inches(0.32), Inches(0.72), SW - Inches(3.2), Inches(0.28))
    _clear(sub_tb.text_frame)
    r2 = sub_tb.text_frame.paragraphs[0].add_run()
    r2.text = f"Synthèse des alertes — {instance_name}  ·  {instance_date}"
    r2.font.size = Pt(8.5); r2.font.name = FONT; r2.font.color.rgb = LIGHT

    # ---- Données ----
    proj_map = {p["project_id"]: p for p in projects}
    red_projects = [p for p in projects if p.get("status_rag") == "red"]
    eac_overruns = [p for p in projects
                    if (p.get("budget_forecast") or 0) > (p.get("budget_total") or 1) * 1.10
                    and (p.get("budget_total") or 0) > 0]
    critical_risks = sorted(
        [r for r in all_risks if r.get("criticality", 0) >= 15 and r.get("status", "") in ("identifié", "en cours")],
        key=lambda x: -x.get("criticality", 0)
    )[:8]

    BODY_Y = HEADER_H + Inches(0.12)
    COL_W = Inches(4.1)
    GUTTER = Inches(0.18)

    def _section(x, y, w, title, icon_color, rows, empty_msg, row_fn):
        """Renders a section box with title + rows."""
        # Section header
        _rect(slide, x, y, w, Inches(0.32), fill=icon_color)
        htb = _tb(slide, x + Inches(0.12), y + Inches(0.04), w - Inches(0.2), Inches(0.24))
        _clear(htb.text_frame)
        rh = htb.text_frame.paragraphs[0].add_run()
        rh.text = title
        rh.font.size = Pt(9); rh.font.bold = True; rh.font.name = FONT; rh.font.color.rgb = WHITE

        # Rows
        ROW_H = Inches(0.37)
        ry = y + Inches(0.32)
        if not rows:
            _rect(slide, x, ry, w, ROW_H, fill=BG, no_line=False, line_color=BORDER, line_pt=0.5)
            mt = _tb(slide, x + Inches(0.12), ry + Inches(0.08), w - Inches(0.2), Inches(0.22))
            _clear(mt.text_frame)
            r = mt.text_frame.paragraphs[0].add_run()
            r.text = empty_msg; r.font.size = Pt(8); r.font.name = FONT; r.font.color.rgb = LIGHT
            ry += ROW_H
        else:
            for item in rows[:6]:
                bg, text_left, text_right = row_fn(item)
                _rect(slide, x, ry, w, ROW_H, fill=bg, no_line=False, line_color=BORDER, line_pt=0.5)
                lt = _tb(slide, x + Inches(0.12), ry + Inches(0.06), w - Inches(1.5), Inches(0.26))
                _clear(lt.text_frame)
                r = lt.text_frame.paragraphs[0].add_run()
                r.text = trunc(text_left, 55); r.font.size = Pt(8); r.font.name = FONT; r.font.color.rgb = DARK
                rt = _tb(slide, x + w - Inches(1.4), ry + Inches(0.06), Inches(1.3), Inches(0.26))
                _clear(rt.text_frame)
                rr = rt.text_frame.paragraphs[0]
                rr.alignment = PP_ALIGN.RIGHT
                rn = rr.add_run()
                rn.text = text_right; rn.font.size = Pt(8); rn.font.bold = True; rn.font.name = FONT; rn.font.color.rgb = DARK
                ry += ROW_H
        return ry

    # Section 1 — Risques critiques non mitigés (left col)
    x1 = Inches(0.18)
    _section(
        x1, BODY_Y, COL_W,
        f"Risques critiques non mitigés ({len(critical_risks)})",
        RED,
        critical_risks,
        "Aucun risque critique non mitigé",
        lambda r: (
            LIGHT_RED,
            f"[{r.get('criticality', 0)}/25] {r.get('title', '?')}",
            proj_map.get(r.get("project_id", ""), {}).get("name", "—")[:22]
        )
    )

    # Section 2 — Projets rouges (middle col)
    x2 = x1 + COL_W + GUTTER
    _section(
        x2, BODY_Y, COL_W,
        f"Projets en statut ROUGE ({len(red_projects)})",
        RED,
        red_projects,
        "Aucun projet rouge",
        lambda p: (
            LIGHT_RED,
            p.get("name", "?"),
            rag_label(p.get("status_rag", ""))
        )
    )

    # Section 3 — EAC dépassés > 10% (right col)
    x3 = x2 + COL_W + GUTTER
    w3 = SW - x3 - Inches(0.18)
    _section(
        x3, BODY_Y, w3,
        f"EAC dépassés > 10 % ({len(eac_overruns)})",
        ORANGE_C,
        eac_overruns,
        "Aucun dépassement EAC > 10 %",
        lambda p: (
            LIGHT_AMBER,
            p.get("name", "?"),
            f"+{round((p.get('budget_forecast', 0) / max(p.get('budget_total', 1), 1) - 1) * 100)} %"
        )
    )

    # Footer
    ftb = _tb(slide, Inches(0.18), SH - Inches(0.35), SW - Inches(0.36), Inches(0.28))
    _clear(ftb.text_frame)
    fp = ftb.text_frame.paragraphs[0]
    fp.alignment = PP_ALIGN.RIGHT
    fr = fp.add_run()
    fr.text = f"{co}  ·  {instance_name}  ·  {instance_date}"
    fr.font.size = Pt(7); fr.font.name = FONT; fr.font.color.rgb = LIGHT


# ====================================================================
# MAIN ENTRY POINT
# ====================================================================

def generate_copil_pptx(instance_name, instance_date, projects,
                        all_milestones, all_risks, all_decisions,
                        governance_id=None, team_consumption_by_project=None,
                        include_roadmap=False, dependencies=None, branding=None):
    global _CURRENT_BRAND
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    brand = _brand(branding)
    _CURRENT_BRAND = brand  # Appliqué à toutes les slides via _run() et _footer()

    proj_names = {p["project_id"]: p["name"] for p in projects}
    tc_by_proj = team_consumption_by_project or {}

    # Enrich decisions with project_name
    for d in all_decisions:
        d["project_name"] = proj_names.get(d.get("project_id", ""), "?")

    # Trier les projets alphabétiquement (spéc : "ordre alphabétique")
    projects_sorted = sorted(projects, key=lambda p: (p.get("name") or "").lower())

    # Slide 1 — Garde
    add_slide_garde(prs, instance_name, instance_date, projects, brand=brand)

    # Slide 2 — Sommaire
    add_slide_sommaire(prs, projects, instance_name, instance_date, brand=brand)

    # Slide 3 — Heatmap
    if all_risks:
        add_slide_heatmap(prs, all_risks, proj_names, instance_name, instance_date, brand=brand)

    # Slide 4 — Top risques
    if all_risks:
        add_slide_top_risks(prs, all_risks, proj_names, instance_name, instance_date, brand=brand)

    # Slide 5 — Décisions
    if all_decisions:
        add_slide_decisions(prs, all_decisions[:10], governance_id, instance_name, instance_date, brand=brand)

    # Slide 6+N — Fiche par projet + slide consommation équipe (ordre alphabétique)
    for p in projects_sorted:
        pid = p["project_id"]
        ms = [m for m in all_milestones if m.get("project_id") == pid]
        ms_sorted = sorted(ms, key=lambda x: x.get("date_forecast") or "")
        r_proj = sorted([r for r in all_risks if r.get("project_id") == pid],
                        key=lambda x: -x.get("criticality", 0))
        d_proj = [d for d in all_decisions if d.get("project_id") == pid]
        add_slide_fiche(prs, p, ms_sorted, r_proj, d_proj, instance_name, instance_date, brand=brand)

        # S1-08 — Slide consommation par équipe
        team_rows = tc_by_proj.get(pid, [])
        add_slide_team_consumption(prs, p, team_rows, instance_name, instance_date, brand=brand)

    # S2-04 — Slide Roadmap consolidée (optionnel)
    if include_roadmap:
        add_slide_roadmap(prs, projects, all_milestones, instance_name, instance_date,
                          dependencies=dependencies, brand=brand)

    # Slide finale — Points d'attention CIO
    add_slide_cloture(prs, projects, all_risks, instance_name, instance_date, brand=brand)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
