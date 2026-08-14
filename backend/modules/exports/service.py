"""Exports PowerPoint — COPIL portefeuille + reporting dédié par instance."""
from datetime import date

from core.database import db
from core.pptx import MarcelDeck, GREEN, AMBER, RED, INDIGO, MUTED, BLUE

RAG_COLOR = {"green": GREEN, "orange": AMBER, "red": RED}
RAG_LABEL = {"green": "Vert", "orange": "Orange", "red": "Rouge"}
SEV_COLOR = {"critical": RED, "critique": RED, "high": RED, "haute": RED,
             "medium": AMBER, "moyenne": AMBER, "low": GREEN, "basse": GREEN}


def _eur(v) -> str:
    v = v or 0
    if abs(v) >= 1_000_000:
        return f"{v/1_000_000:.1f} M€".replace(".", ",")
    if abs(v) >= 1_000:
        return f"{v/1_000:.0f} K€"
    return f"{v:.0f} €"


async def _deck(user) -> MarcelDeck:
    tenant = await db.tenants.find_one({"tenant_id": user.tenant_id}, {"_id": 0, "name": 1}) or {}
    return MarcelDeck(tenant.get("name", ""))


# ─── Slides réutilisables ──────────────────────────────────────────────────────

async def _copil_slides(deck, user):
    projects = await db.projects.find(
        {"tenant_id": user.tenant_id},
        {"_id": 0, "project_id": 1, "name": 1, "code": 1, "status": 1, "status_rag": 1,
         "budget_total": 1, "budget_consumed": 1, "eac": 1, "end_date": 1}).to_list(None)
    actifs = [p for p in projects if p.get("status") == "actif"]
    rag = {"green": 0, "orange": 0, "red": 0}
    for p in actifs:
        rag[p.get("status_rag", "green")] = rag.get(p.get("status_rag", "green"), 0) + 1
    budget = sum(p.get("budget_total") or 0 for p in actifs)
    consumed = sum(p.get("budget_consumed") or 0 for p in actifs)

    deck.kpis("Vue d'ensemble du portefeuille", [
        {"label": "Projets actifs", "value": len(actifs)},
        {"label": "Verts", "value": rag["green"], "color": GREEN},
        {"label": "Orange", "value": rag["orange"], "color": AMBER},
        {"label": "Rouges", "value": rag["red"], "color": RED},
        {"label": "Budget engagé", "value": _eur(budget), "hint": f"Consommé : {_eur(consumed)}"},
    ])

    top = sorted(actifs, key=lambda p: -(p.get("budget_total") or 0))[:10]
    rows, colors = [], {}
    for i, p in enumerate(top):
        r = p.get("status_rag", "green")
        rows.append([p.get("code") or "—", p.get("name", "")[:52], RAG_LABEL.get(r, r),
                     _eur(p.get("budget_total")), _eur(p.get("budget_consumed")), _eur(p.get("eac"))])
        colors[(i, 2)] = RAG_COLOR.get(r, INDIGO)
    deck.table("Projets du portefeuille", ["Code", "Projet", "RAG", "Budget", "Consommé", "EAC"],
               rows, subtitle="Top 10 par budget", col_widths=[1.1, 4.2, 0.9, 1.3, 1.3, 1.3], cell_colors=colors)

    alerts = [p for p in actifs if p.get("status_rag") in ("red", "orange")]
    if alerts:
        arows, acolors = [], {}
        for i, p in enumerate(sorted(alerts, key=lambda x: 0 if x.get("status_rag") == "red" else 1)[:10]):
            r = p.get("status_rag")
            over = (p.get("eac") or 0) - (p.get("budget_total") or 0)
            arows.append([p.get("code") or "—", p.get("name", "")[:50], RAG_LABEL.get(r, r),
                          _eur(over) if over > 0 else "—", p.get("end_date") or "—"])
            acolors[(i, 2)] = RAG_COLOR.get(r, INDIGO)
        deck.table("Projets sous surveillance", ["Code", "Projet", "RAG", "Dépassement EAC", "Fin prévue"],
                   arows, col_widths=[1.1, 4.5, 0.9, 1.6, 1.4], cell_colors=acolors)

    # Écart budget / forecast (reforecast trimestriel)
    from modules.forecast.service import quarters_summary
    fc = await quarters_summary(date.today().year, user.tenant_id)
    ecart = (fc["totals"]["budget"] or 0) - (fc["totals"]["forecast"] or 0)
    deck.kpis(f"Reforecast {date.today().year} — écart budget / forecast", [
        {"label": "Budget portefeuille", "value": _eur(fc["totals"]["budget"])},
        {"label": "Forecast (scope valorisé)", "value": _eur(fc["totals"]["forecast"]), "color": BLUE},
        {"label": "Consommé", "value": _eur(fc["totals"]["consumed"])},
        {"label": "Écart", "value": _eur(ecart), "color": RED if ecart < 0 else GREEN},
    ], subtitle="Forecast = JH alloués × TJM réel des ressources")

    decisions = await db.decisions.find(
        {"tenant_id": user.tenant_id}, {"_id": 0, "title": 1, "status": 1, "decision_date": 1, "owner": 1}
    ).sort("decision_date", -1).limit(8).to_list(None)
    if decisions:
        deck.table("Décisions récentes", ["Décision", "Statut", "Date", "Porteur"],
                   [[d.get("title", "")[:60], d.get("status", ""), d.get("decision_date") or "—",
                     d.get("owner") or "—"] for d in decisions],
                   col_widths=[4.6, 1.2, 1.2, 1.6])

    events = await db.events.find(
        {"tenant_id": user.tenant_id, "date": {"$gte": date.today().isoformat()}, "status": "planifie"},
        {"_id": 0, "title": 1, "date": 1, "level": 1}).sort("date", 1).limit(10).to_list(None)
    if events:
        deck.table("Prochaines instances", ["Date", "Instance", "Niveau"],
                   [[e["date"], e["title"], e["level"].capitalize()] for e in events],
                   subtitle="Calendrier des instances MARCEL", col_widths=[1.2, 4.5, 1.5])


async def _reforecast_slides(deck, user):
    from modules.forecast.service import quarters_summary
    year = date.today().year
    fc = await quarters_summary(year, user.tenant_id)
    ecart = (fc["totals"]["budget"] or 0) - (fc["totals"]["forecast"] or 0)
    deck.kpis(f"Reforecast {year} — synthèse", [
        {"label": "Budget", "value": _eur(fc["totals"]["budget"])},
        {"label": "Forecast", "value": _eur(fc["totals"]["forecast"]), "color": BLUE},
        {"label": "Consommé", "value": _eur(fc["totals"]["consumed"])},
        {"label": "Écart", "value": _eur(ecart), "color": RED if ecart < 0 else GREEN},
        {"label": "TJM moyen", "value": f"{fc['default_tjm']} €"},
    ])
    rows, colors = [], {}
    for i, p in enumerate(fc["projects"][:12]):
        cells = [(_eur(c["final_value"] if c["validated"] else c["scope_value"]) + (" ✓" if c["validated"] else ""))
                 for c in p["quarters"]]
        rows.append([p.get("code") or "—", p.get("name", "")[:36], *cells,
                     _eur(p["forecast_year"]), _eur(p["ecart_budget"])])
        colors[(i, 7)] = RED if p["ecart_budget"] < 0 else GREEN
    deck.table(f"Reforecast trimestriel {year}", ["Code", "Projet", "Q1", "Q2", "Q3", "Q4", "Forecast", "Écart"],
               rows, subtitle="Scope valorisé en € — ✓ = trimestre validé",
               col_widths=[0.9, 3.0, 1.0, 1.0, 1.0, 1.0, 1.2, 1.2], cell_colors=colors)
    transfers = await db.budget_transfers.find(
        {"tenant_id": user.tenant_id}, {"_id": 0}).sort("created_at", -1).limit(8).to_list(None)
    if transfers:
        deck.table("Transferts budgétaires récents", ["Date", "De", "Vers", "Montant", "Motif"],
                   [[t["created_at"][:10], t.get("from_project_name", "")[:26], t.get("to_project_name", "")[:26],
                     _eur(t["amount"]), (t.get("reason") or "—")[:38]] for t in transfers],
                   col_widths=[1.0, 2.2, 2.2, 1.1, 3.0])


async def _gates_slides(deck, user):
    gates = await db.lifecycle_gates.find(
        {"tenant_id": user.tenant_id}, {"_id": 0}).sort("target_date", 1).to_list(None)
    pending = [g for g in gates if g.get("status") in ("en_validation", "pret", "demande")]
    decided = [g for g in gates if g.get("status") in ("go", "no_go", "go_reserves")][-8:]
    deck.kpis("Comité d'investissement — gates", [
        {"label": "Passages en cours", "value": len(pending), "color": BLUE},
        {"label": "Prêts pour décision", "value": len([g for g in pending if g.get("status") == "pret"]), "color": GREEN},
        {"label": "Décisions rendues", "value": len([g for g in gates if g.get("status") in ("go", "no_go", "go_reserves")])},
    ])
    if pending:
        deck.table("Demandes de passage de phase", ["Projet", "Passage", "Statut", "Date cible", "Livrables OK"],
                   [[f"{g.get('project_code') or ''} {g.get('project_name', '')}"[:38],
                     f"{g.get('from_phase')} → {g.get('to_phase')}", g.get("status", ""),
                     g.get("target_date") or "—",
                     f"{sum(1 for d in g.get('deliverables', []) if d.get('review_status') == 'valide')}/{len(g.get('deliverables', []))}"]
                    for g in pending[:10]],
                   col_widths=[3.2, 1.8, 1.2, 1.2, 1.2])
        g = pending[0]
        if g.get("deliverables"):
            rows, colors = [], {}
            for i, d in enumerate(g["deliverables"]):
                st = d.get("review_status", "pending")
                rows.append([d.get("label", d.get("key", "")), d.get("validator") or "—",
                             "Fourni" if d.get("provided") else "Attendu",
                             {"valide": "Validé", "refuse": "Refusé", "pending": "En attente"}.get(st, st),
                             d.get("reviewed_by_name") or "—"])
                colors[(i, 3)] = GREEN if st == "valide" else RED if st == "refuse" else MUTED
            deck.table(f"Dossier de gate — {g.get('project_name', '')[:40]}",
                       ["Livrable", "Validateur", "Fourniture", "Avis", "Par"], rows,
                       subtitle=f"Passage {g.get('from_phase')} → {g.get('to_phase')}",
                       col_widths=[3.0, 1.4, 1.2, 1.2, 1.8], cell_colors=colors)
    if decided:
        rows, colors = [], {}
        for i, g in enumerate(decided):
            dec = g.get("status", "")
            rows.append([f"{g.get('project_code') or ''} {g.get('project_name', '')}"[:40],
                         f"{g.get('from_phase')} → {g.get('to_phase')}",
                         {"go": "GO", "no_go": "NO-GO", "go_reserves": "GO avec réserves"}.get(dec, dec),
                         (g.get("decision") or {}).get("decided_at", "")[:10] if isinstance(g.get("decision"), dict) else "—"])
            colors[(i, 2)] = GREEN if dec == "go" else RED if dec == "no_go" else AMBER
        deck.table("Décisions de gate récentes", ["Projet", "Passage", "Décision", "Date"], rows,
                   col_widths=[3.6, 1.8, 1.6, 1.2], cell_colors=colors)


async def _strategic_slides(deck, user):
    from modules.budget_ops.service import list_envelopes
    year = date.today().year
    data = await list_envelopes(user.tenant_id, year)
    envs = data["envelopes"]
    total = sum(e.get("amount") or 0 for e in envs)
    engaged = sum(e.get("engaged") or 0 for e in envs)
    deck.kpis(f"Enveloppes stratégiques {year}", [
        {"label": "Enveloppes définies", "value": len(envs)},
        {"label": "Montant total", "value": _eur(total), "color": INDIGO},
        {"label": "Engagé (projets)", "value": _eur(engaged), "color": BLUE},
        {"label": "Taux d'engagement", "value": f"{round(engaged / total * 100) if total else 0} %"},
    ])
    for axis, label in (("programme", "Par programme"), ("theme", "Par thème stratégique")):
        sub = [e for e in envs if e["axis"] == axis]
        if sub:
            rows, colors = [], {}
            for i, e in enumerate(sub):
                rows.append([e["ref_name"][:40], _eur(e["amount"]), _eur(e["engaged"]), _eur(e["consumed"]), f"{e['rate']} %"])
                colors[(i, 4)] = RED if e["rate"] > 100 else AMBER if e["rate"] >= 85 else GREEN
            deck.table(f"Enveloppes {label.lower()}", ["Axe", "Enveloppe", "Engagé", "Consommé", "Taux"],
                       rows, col_widths=[3.4, 1.4, 1.4, 1.4, 0.9], cell_colors=colors)
    if not envs:
        deck.bullets("Enveloppes stratégiques", ["Aucune enveloppe définie pour l'année — à valider en comité stratégique (Budget → Enveloppes)."])


async def _security_slides(deck, user):
    vulns = await db.vulnerabilities.find({"tenant_id": user.tenant_id}, {"_id": 0}).to_list(None)
    open_v = [v for v in vulns if v.get("status") not in ("resolved", "corrigee", "closed", "fermee")]
    counts = {}
    for v in open_v:
        counts[v.get("severity", "?")] = counts.get(v.get("severity", "?"), 0) + 1
    deck.kpis("Revue sécurité — vulnérabilités ouvertes", [
        {"label": "Total ouvertes", "value": len(open_v)},
        {"label": "Critiques", "value": counts.get("critical", 0) + counts.get("critique", 0), "color": RED},
        {"label": "Hautes", "value": counts.get("high", 0) + counts.get("haute", 0), "color": RED},
        {"label": "Moyennes", "value": counts.get("medium", 0) + counts.get("moyenne", 0), "color": AMBER},
        {"label": "Basses", "value": counts.get("low", 0) + counts.get("basse", 0), "color": GREEN},
    ])
    if open_v:
        order = {"critical": 0, "critique": 0, "high": 1, "haute": 1, "medium": 2, "moyenne": 2}
        top = sorted(open_v, key=lambda v: order.get(v.get("severity"), 3))[:12]
        rows, colors = [], {}
        for i, v in enumerate(top):
            rows.append([v.get("title", "")[:52], v.get("severity", ""), v.get("source") or "—",
                         v.get("status", ""), v.get("due_date") or "—"])
            colors[(i, 1)] = SEV_COLOR.get(v.get("severity"), MUTED)
        deck.table("Plan de remédiation", ["Vulnérabilité", "Sévérité", "Source", "Statut", "Échéance"],
                   rows, col_widths=[3.8, 1.1, 1.2, 1.1, 1.2], cell_colors=colors)
    else:
        deck.bullets("Vulnérabilités", ["Aucune vulnérabilité ouverte — RAS."])


async def _mep_slides(deck, user):
    today = date.today().isoformat()
    releases = await db.releases.find({"tenant_id": user.tenant_id}, {"_id": 0}).sort("date", 1).to_list(None)
    upcoming = [r for r in releases if (r.get("date") or "") >= today][:10]
    recent = [r for r in releases if (r.get("date") or "") < today][-8:]
    deck.kpis("CAB — mises en production", [
        {"label": "MEP à venir", "value": len(upcoming), "color": BLUE},
        {"label": "MEP réalisées", "value": len(recent)},
    ])
    if upcoming:
        deck.table("MEP à valider / à venir", ["Date", "Release", "Type", "Statut"],
                   [[r.get("date") or "—", r.get("name", "")[:48], r.get("type") or "—", r.get("status", "")]
                    for r in upcoming], col_widths=[1.2, 4.2, 1.2, 1.4])
    if recent:
        deck.table("MEP récentes", ["Date", "Release", "Type", "Statut"],
                   [[r.get("date") or "—", r.get("name", "")[:48], r.get("type") or "—", r.get("status", "")]
                    for r in recent], col_widths=[1.2, 4.2, 1.2, 1.4])
    if not releases:
        deck.bullets("Mises en production", ["Aucune MEP enregistrée dans le module Run."])


async def _capacity_slides(deck, user):
    from modules.capacity.service import console
    data = await console(user.tenant_id, 3, "team")
    months = [m[:7] for m in data["months"]]
    over = [r for r in data["rows"] if r["rate"] > 100]
    deck.kpis("Revue de capacité — 3 mois", [
        {"label": "Capacité", "value": f"{data['totals']['capacity']} JH"},
        {"label": "Charge", "value": f"{data['totals']['load']} JH", "color": BLUE},
        {"label": "Équipes en surcharge", "value": len(over), "color": RED if over else GREEN},
    ])
    rows, colors = [], {}
    for i, r in enumerate(data["rows"][:12]):
        cells = [f"{c['rate']} %" if c["capacity"] else "—" for c in r["cells"]]
        rows.append([r["label"][:30], str(r["resources"]), *cells, f"{r['rate']} %"])
        colors[(i, 5)] = RED if r["rate"] > 100 else AMBER if r["rate"] >= 85 else GREEN
    deck.table("Charge vs capacité par équipe", ["Équipe", "Res.", *months, "Total"],
               rows, subtitle="Taux = JH alloués / capacité disponible",
               col_widths=[2.6, 0.7, 1.2, 1.2, 1.2, 1.0], cell_colors=colors)


async def _vendors_slides(deck, user):
    externals = await db.resources.find(
        {"tenant_id": user.tenant_id, "resource_type": {"$in": ["externe_regie", "externe_forfait"]}},
        {"_id": 0, "name": 1, "vendor": 1, "contract_tjm": 1, "contract_start": 1,
         "contract_end": 1, "contract_ref": 1, "resource_type": 1}).to_list(None)
    today = date.today().isoformat()
    deck.kpis("Revue fournisseurs", [
        {"label": "Externes actifs", "value": len(externals)},
        {"label": "ESN distinctes", "value": len({r.get("vendor") for r in externals if r.get("vendor")})},
        {"label": "Contrats échus", "value": len([r for r in externals if (r.get("contract_end") or "9999") < today]), "color": RED},
    ])
    if externals:
        rows, colors = [], {}
        ordered = sorted(externals, key=lambda r: r.get("contract_end") or "9999")[:12]
        for i, r in enumerate(ordered):
            expired = (r.get("contract_end") or "9999") < today
            rows.append([r.get("name", ""), r.get("vendor") or "—", r.get("resource_type", "").replace("externe_", ""),
                         f"{r.get('contract_tjm') or '—'} €", r.get("contract_end") or "—", r.get("contract_ref") or "—"])
            colors[(i, 4)] = RED if expired else INDIGO
        deck.table("Contrats externes", ["Ressource", "ESN", "Type", "TJM", "Fin contrat", "Réf."],
                   rows, subtitle="Triés par échéance de contrat", col_widths=[2.2, 1.8, 1.0, 0.9, 1.2, 1.4], cell_colors=colors)
    else:
        deck.bullets("Fournisseurs", ["Aucune ressource externe enregistrée."])


async def _safe_slides(deck, user):
    trains = await db.trains.find({"tenant_id": user.tenant_id}, {"_id": 0}).to_list(None)
    pis = await db.pis.find({"tenant_id": user.tenant_id}, {"_id": 0}).sort("start_date", -1).to_list(None)
    deck.kpis("SAFe — trains & PIs", [
        {"label": "Trains (ART)", "value": len(trains)},
        {"label": "PIs", "value": len(pis)},
        {"label": "PIs en cours", "value": len([p for p in pis if p.get("status") in ("active", "en_cours", "in_progress")]), "color": BLUE},
    ])
    if pis:
        tnames = {t["train_id"]: t["name"] for t in trains}
        deck.table("Program Increments", ["PI", "ART", "Début", "Fin", "Objectifs", "Statut"],
                   [[p.get("name", ""), tnames.get(p.get("train_id"), "—")[:22], p.get("start_date") or "—",
                     p.get("end_date") or "—", str(len(p.get("objectives") or [])), p.get("status", "")]
                    for p in pis[:10]], col_widths=[1.6, 2.0, 1.1, 1.1, 1.0, 1.2])
        latest = pis[0]
        if latest.get("objectives"):
            deck.bullets(f"Objectifs du PI — {latest.get('name', '')}", latest["objectives"][:10])
    if not trains:
        deck.bullets("SAFe", ["Aucun train configuré."])


async def _trajectory_slides(deck, user):
    apps = await db.applications.find(
        {"tenant_id": user.tenant_id},
        {"_id": 0, "name": 1, "criticality": 1, "disposition": 1, "trajectory_target_date": 1}).to_list(None)
    ms = await db.trajectory_milestones.find({"tenant_id": user.tenant_id}, {"_id": 0}).sort("date", 1).to_list(None)
    counts = {}
    for a in apps:
        counts[a.get("disposition") or "non_positionne"] = counts.get(a.get("disposition") or "non_positionne", 0) + 1
    deck.kpis("Trajectoire du SI (TIME)", [
        {"label": "Conserver", "value": counts.get("conserver", 0), "color": GREEN},
        {"label": "Moderniser", "value": counts.get("moderniser", 0), "color": BLUE},
        {"label": "Remplacer", "value": counts.get("remplacer", 0), "color": AMBER},
        {"label": "Décommissionner", "value": counts.get("decommissionner", 0), "color": RED},
        {"label": "Non positionnées", "value": counts.get("non_positionne", 0)},
    ])
    positioned = [a for a in apps if a.get("disposition")]
    if positioned:
        deck.table("Dispositions applicatives", ["Application", "Criticité", "Disposition", "Cible"],
                   [[a.get("name", "")[:40], a.get("criticality") or "—", a.get("disposition", ""),
                     a.get("trajectory_target_date") or "—"] for a in positioned[:14]],
                   col_widths=[3.4, 1.2, 1.5, 1.2])
    if ms:
        deck.table("Jalons de trajectoire", ["Date", "Jalon", "Statut"],
                   [[m.get("date") or "—", m.get("title", "")[:56], "Fait" if m.get("status") == "fait" else "À venir"]
                    for m in ms[:12]], col_widths=[1.2, 4.6, 1.1])


async def _demands_slides(deck, user):
    demands = await db.demands.find({"tenant_id": user.tenant_id}, {"_id": 0}).sort("created_at", -1).to_list(None)
    pending = [d for d in demands if d.get("status") in ("soumise", "submitted", "nouvelle", "new", "a_qualifier", "pending")]
    deck.kpis("Comité des demandes", [
        {"label": "Demandes totales", "value": len(demands)},
        {"label": "À qualifier", "value": len(pending), "color": AMBER},
        {"label": "Budget estimé (à qualifier)", "value": _eur(sum(d.get("estimated_budget") or 0 for d in pending))},
    ])
    show = pending or demands
    if show:
        deck.table("Demandes", ["Demande", "Demandeur", "Direction", "Budget est.", "Urgence", "Statut", "Score"],
                   [[d.get("title", "")[:36], d.get("requester") or "—", (d.get("requester_department") or "—")[:16],
                     _eur(d.get("estimated_budget")), str(d.get("urgency") or "—"), d.get("status", ""),
                     str(d.get("priority_score") or "—")] for d in show[:12]],
                   col_widths=[2.6, 1.4, 1.3, 1.0, 0.8, 1.1, 0.7])
    else:
        deck.bullets("Demandes", ["Aucune demande enregistrée."])


async def _projects_risks_slides(deck, user):
    projects = await db.projects.find(
        {"tenant_id": user.tenant_id, "status": "actif"},
        {"_id": 0, "project_id": 1, "name": 1, "code": 1, "status_rag": 1, "phase": 1, "end_date": 1}).to_list(None)
    pmap = {p["project_id"]: p for p in projects}
    risks = await db.risks.find(
        {"tenant_id": user.tenant_id, "status": {"$nin": ["closed", "ferme", "clos"]}}, {"_id": 0}).to_list(None)
    risks.sort(key=lambda r: -(r.get("criticality") or 0))
    today = date.today().isoformat()
    milestones = await db.milestones.find(
        {"tenant_id": user.tenant_id, "status": {"$ne": "done"},
         "date_forecast": {"$gte": today}}, {"_id": 0}).sort("date_forecast", 1).limit(10).to_list(None)

    rows, colors = [], {}
    for i, p in enumerate(sorted(projects, key=lambda x: 0 if x.get("status_rag") == "red" else 1 if x.get("status_rag") == "orange" else 2)[:12]):
        r = p.get("status_rag", "green")
        rows.append([p.get("code") or "—", p.get("name", "")[:44], (p.get("phase") or "—").capitalize(),
                     RAG_LABEL.get(r, r), p.get("end_date") or "—"])
        colors[(i, 3)] = RAG_COLOR.get(r, INDIGO)
    deck.table("Avancement des projets actifs", ["Code", "Projet", "Phase", "RAG", "Fin prévue"],
               rows, col_widths=[0.9, 3.8, 1.3, 0.9, 1.2], cell_colors=colors)

    if risks:
        rrows, rcolors = [], {}
        for i, r in enumerate(risks[:10]):
            crit = r.get("criticality") or 0
            rrows.append([r.get("title", "")[:42], pmap.get(r.get("project_id"), {}).get("name", "—")[:26],
                          str(crit), r.get("owner") or "—", (r.get("mitigation_plan") or "—")[:30]])
            rcolors[(i, 2)] = RED if crit >= 15 else AMBER if crit >= 8 else GREEN
        deck.table("Top risques ouverts", ["Risque", "Projet", "Criticité", "Porteur", "Plan d'action"],
                   rrows, col_widths=[2.8, 1.9, 0.9, 1.3, 2.0], cell_colors=rcolors)

    if milestones:
        deck.table("Prochains jalons", ["Date", "Jalon", "Projet", "Statut"],
                   [[m.get("date_forecast") or "—", m.get("name", "")[:40],
                     pmap.get(m.get("project_id"), {}).get("name", "—")[:26], m.get("status", "")]
                    for m in milestones], col_widths=[1.1, 3.2, 2.2, 1.1])


# ─── Mapping type d'instance → builder ─────────────────────────────────────────

BUILDERS = [
    (("reforecast",), _reforecast_slides),
    (("investissement", "gate"), _gates_slides),
    (("stratégique", "strategique", "cadrage"), _strategic_slides),
    (("sécurité", "securite"), _security_slides),
    (("cab", "mep"), _mep_slides),
    (("capacité", "capacite"), _capacity_slides),
    (("fournisseur",), _vendors_slides),
    (("pi planning", "inspect", "system demo", "portfolio sync", "participatory"), _safe_slides),
    (("trajectoire", "applicatif", "apm"), _trajectory_slides),
    (("demandes",), _demands_slides),
    (("copil projet", "coproj", "avancement", "risques"), _projects_risks_slides),
]


def _resolve_builder(title: str):
    t = (title or "").lower()
    for keywords, fn in BUILDERS:
        if any(k in t for k in keywords):
            return fn
    return _copil_slides


async def _roadmap_slides(deck, user):
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE
    from core.pptx import _round_rect, _text, _set_fill, CARD, BG, TEXT, MUTED, BORDER
    from pptx.dml.color import RGBColor
    from datetime import timedelta

    today = date.today()
    win_start = date(today.year, today.month, 1)
    m = win_start.month - 3
    y = win_start.year
    while m < 1:
        m += 12
        y -= 1
    win_start = date(y, m, 1)
    months = []
    my, mm = win_start.year, win_start.month
    for _ in range(12):
        months.append(date(my, mm, 1))
        mm += 1
        if mm > 12:
            mm, my = 1, my + 1
    end_y, end_m = months[-1].year, months[-1].month + 1
    if end_m > 12:
        end_m, end_y = 1, end_y + 1
    win_end = date(end_y, end_m, 1)
    span = (win_end - win_start).days

    projects = await db.projects.find(
        {"tenant_id": user.tenant_id, "status": {"$nin": ["annule"]}},
        {"_id": 0, "project_id": 1, "name": 1, "code": 1, "status_rag": 1, "program_id": 1,
         "start_date": 1, "end_date_forecast": 1, "end_date_baseline": 1}).to_list(None)
    programs = await db.programs.find({"tenant_id": user.tenant_id}, {"_id": 0, "program_id": 1, "name": 1}).to_list(None)
    pnames = {p["program_id"]: p["name"] for p in programs}
    milestones = await db.milestones.find(
        {"tenant_id": user.tenant_id, "is_governance": True}, {"_id": 0, "project_id": 1, "date_forecast": 1, "date_baseline": 1}).to_list(None)
    ms_by_p: dict = {}
    for msn in milestones:
        ms_by_p.setdefault(msn["project_id"], []).append(msn.get("date_forecast") or msn.get("date_baseline"))

    groups: dict = {}
    for p in projects:
        groups.setdefault(p.get("program_id") or "_none", []).append(p)
    rows = []
    for gid in sorted(groups, key=lambda g: pnames.get(g, "zzz")):
        rows.append(("group", pnames.get(gid, "Sans programme")))
        for p in sorted(groups[gid], key=lambda x: x.get("start_date") or ""):
            rows.append(("project", p))

    GRID_X, GRID_W = Inches(3.0), Inches(9.7)
    ROW_H, START_Y = Inches(0.34), Inches(1.85)
    PER_SLIDE = 14

    def x_of(dstr):
        try:
            d = date.fromisoformat(dstr[:10])
        except (ValueError, TypeError):
            return None
        frac = (d - win_start).days / span
        return max(0.0, min(frac, 1.0))

    for chunk_start in range(0, len(rows), PER_SLIDE):
        chunk = rows[chunk_start:chunk_start + PER_SLIDE]
        slide = deck._new()
        deck._header(slide, "Roadmap Portefeuille",
                     f"{win_start.strftime('%m/%Y')} → {months[-1].strftime('%m/%Y')} — barres RAG, losanges = jalons de gouvernance")
        col_w = Emu(int(GRID_W / 12))
        for i, mo in enumerate(months):
            _text(slide, Emu(int(GRID_X + col_w * i)), Inches(1.5), col_w, Inches(0.3),
                  mo.strftime("%m/%y"), size=8, color=MUTED, bold=True)
            sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(GRID_X + col_w * i)), START_Y,
                                         Emu(6350), Emu(int(ROW_H * len(chunk))))
            _set_fill(sep, BORDER)
            sep.shadow.inherit = False
        tx = x_of(today.isoformat())
        if tx is not None:
            tl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(GRID_X + GRID_W * tx)), Inches(1.5),
                                        Emu(19050), Emu(int(Inches(0.35) + ROW_H * len(chunk))))
            _set_fill(tl, RGBColor(0x2E, 0x5F, 0xE8))
            tl.shadow.inherit = False
        y = START_Y
        for kind, item in chunk:
            if kind == "group":
                _text(slide, Inches(0.6), y, Inches(2.3), ROW_H, str(item)[:34].upper(), size=9, bold=True, color=INDIGO)
            else:
                p = item
                _text(slide, Inches(0.7), y, Inches(2.25), ROW_H,
                      f"{p.get('code') or ''} {p.get('name', '')}"[:38], size=8.5, color=TEXT)
                x1 = x_of(p.get("start_date") or "")
                x2 = x_of(p.get("end_date_forecast") or p.get("end_date_baseline") or "")
                if x1 is not None and x2 is not None and x2 > x1:
                    bar = _round_rect(slide, Emu(int(GRID_X + GRID_W * x1)), Emu(int(y + Inches(0.06))),
                                      Emu(max(int(GRID_W * (x2 - x1)), 50000)), Inches(0.20),
                                      RAG_COLOR.get(p.get("status_rag"), GREEN), radius=0.5)
                for md in ms_by_p.get(p["project_id"], []):
                    mx = x_of(md or "")
                    if mx is None:
                        continue
                    dia = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                                 Emu(int(GRID_X + GRID_W * mx - 45000)), Emu(int(y + Inches(0.05))),
                                                 Emu(90000), Emu(90000))
                    _set_fill(dia, INDIGO)
                    dia.shadow.inherit = False
            y = Emu(int(y + ROW_H))
        deck._footer(slide)


async def build_roadmap_pptx(user) -> bytes:
    deck = await _deck(user)
    deck.cover("Roadmap Portefeuille", "Vue consolidée multi-projets — 12 mois glissants")
    await _roadmap_slides(deck, user)
    deck.section("Merci", "Généré automatiquement par MARCEL")
    return deck.to_bytes()


async def build_copil_pptx(user) -> bytes:
    deck = await _deck(user)
    deck.cover("COPIL Portefeuille", "Comité de pilotage — synthèse mensuelle")
    await _copil_slides(deck, user)
    deck.section("Merci", "Généré automatiquement par MARCEL")
    return deck.to_bytes()


async def build_event_pptx(event_id: str, user):
    event = await db.events.find_one({"event_id": event_id, "tenant_id": user.tenant_id}, {"_id": 0})
    if not event:
        return None, None
    etype = await db.event_types.find_one(
        {"event_type_id": event.get("event_type_id"), "tenant_id": user.tenant_id}, {"_id": 0}) or {}
    deck = await _deck(user)
    d = event.get("date", "")
    pretty = f"{d[8:10]}/{d[5:7]}/{d[:4]}" if len(d) == 10 else d
    deck.cover(event.get("title", "Instance"),
               etype.get("description", "") or f"Instance de niveau {event.get('level', '')}",
               meta=f"{deck.tenant} — séance du {pretty}")
    builder = _resolve_builder(event.get("title", ""))
    await builder(deck, user)
    deck.section("Merci", f"Reporting généré automatiquement par MARCEL — {event.get('title', '')}")
    slug = "".join(c if c.isalnum() else "_" for c in event.get("title", "instance"))[:40]
    return deck.to_bytes(), f"{slug}_{d}.pptx"
