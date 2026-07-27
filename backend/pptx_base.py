"""Générateur PowerPoint COPIL — MARCEL (v2)"""
import io
import base64
from datetime import datetime
from typing import List, Optional

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---- Slide dimensions (widescreen 16:9) ----
SW = Inches(13.33)
SH = Inches(7.5)

# ---- Palette — fond blanc, accents Navy #0B2545, RAG couleurs corrigées ----
NAVY        = RGBColor(0x0B, 0x25, 0x45)
BLUE        = RGBColor(0x00, 0x52, 0xCC)
LIGHT_BLUE  = RGBColor(0xEB, 0xF2, 0xFF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG          = RGBColor(0xF8, 0xF9, 0xFA)
DARK        = RGBColor(0x1E, 0x29, 0x3B)
MID         = RGBColor(0x64, 0x74, 0x8B)
LIGHT       = RGBColor(0x94, 0xA3, 0xB8)
BORDER      = RGBColor(0xE2, 0xE8, 0xF0)
RED         = RGBColor(0xEF, 0x44, 0x44)
ORANGE_C    = RGBColor(0xF5, 0x9E, 0x0B)
GREEN_C     = RGBColor(0x10, 0xB9, 0x81)
LIGHT_RED   = RGBColor(0xFE, 0xF2, 0xF2)
LIGHT_AMBER = RGBColor(0xFF, 0xF7, 0xED)
LIGHT_GREEN = RGBColor(0xF0, 0xFD, 0xF4)

DECISION_STATUS_COLORS = {
    "proposée":  RGBColor(0xE6, 0xF0, 0xFF),
    "prise":     RGBColor(0xED, 0xE9, 0xFE),
    "en_cours":  RGBColor(0xFE, 0xF3, 0xC7),
    "appliquée": RGBColor(0xDC, 0xFC, 0xE7),
    "reportée":  RGBColor(0xF1, 0xF5, 0xF9),
    "annulée":   RGBColor(0xFE, 0xE2, 0xE2),
}

FONT = "Arial"

# ── Branding dynamique (mis à jour à chaque génération) ──────────────────────
_CURRENT_BRAND: dict = {}


# ---- Branding helpers ----

def _parse_rgb(hex_str, default: RGBColor) -> RGBColor:
    """Parse une couleur hex (#RRGGBB) en RGBColor avec fallback."""
    if not hex_str:
        return default
    try:
        h = str(hex_str).lstrip('#')
        if len(h) == 6:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        pass
    return default


def _brand(branding: dict | None) -> dict:
    """Construit le dictionnaire de couleurs de marque depuis la config tenant."""
    b = branding or {}
    return {
        "primary":      _parse_rgb(b.get("primary_color"), NAVY),
        "secondary":    _parse_rgb(b.get("secondary_color"), BLUE),
        "accent":       _parse_rgb(b.get("accent_color"), GREEN_C),
        "company_name": b.get("company_name") or "MARCEL",
        "font":         b.get("font") or FONT,
        "logo_base64":  b.get("logo_base64"),
    }



def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, left, top, width, height, fill=None, no_line=True, line_color=None, line_pt=0.5):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if no_line:
        shape.line.width = 0
    elif line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    return shape


def _tb(slide, left, top, width, height, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = wrap
    return tb


def _clear(tf):
    for p in tf.paragraphs:
        p.clear()


def _run(tf, text, size=9, bold=False, color=None, align=PP_ALIGN.LEFT, space_before=0, italic=False):
    """Add a paragraph with a single run to a text frame (brand font)."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = _CURRENT_BRAND.get("font", FONT) if _CURRENT_BRAND else FONT
    r.font.color.rgb = color or DARK
    return p


def keur(val):
    if val is None:
        return "—"
    return f"{int(val / 1000):,}".replace(",", "\u202f") + " K€"


def eur(val):
    """Format as xxx K€ or xxx € depending on magnitude."""
    if val is None:
        return "—"
    if abs(val) >= 1000:
        return keur(val)
    return f"{int(val):,}".replace(",", "\u202f") + " €"


def jh_fmt(val):
    if val is None:
        return "—"
    return f"{val:.1f} JH" if val != int(val) else f"{int(val)} JH"


def fmt_date(d):
    if not d:
        return "—"
    try:
        return datetime.fromisoformat(str(d).replace("Z", "")).strftime("%d/%m/%Y")
    except Exception:
        return str(d)[:10] if d else "—"


def trunc(s, n=55):
    if not s:
        return "—"
    return s[:n] + "…" if len(s) > n else s


def crit_color(c):
    return RED if c >= 16 else ORANGE_C if c >= 7 else GREEN_C


def crit_bg(c):
    return LIGHT_RED if c >= 16 else LIGHT_AMBER if c >= 7 else LIGHT_GREEN


def rag_color(rag):
    return {"green": GREEN_C, "orange": ORANGE_C, "red": RED}.get(str(rag), MID)


def rag_label(rag):
    return {"green": "Vert", "orange": "Orange", "red": "Rouge"}.get(str(rag), "—")


def status_label(s):
    return {
        "en_preparation": "En préparation", "actif": "Actif",
        "en_pause": "En pause", "cloture": "Clôturé", "archive": "Archivé",
    }.get(str(s), str(s) if s else "—")


def decision_status_label(s):
    return {
        "proposée": "Proposée", "prise": "Prise", "en_cours": "En cours",
        "appliquée": "Appliquée", "reportée": "Reportée", "annulée": "Annulée",
    }.get(str(s), str(s) if s else "—")


def _set_brand(brand: dict) -> None:
    """Met à jour le branding courant in-place (visible par tous les modules importants)."""
    _CURRENT_BRAND.clear()
    _CURRENT_BRAND.update(brand)


# ---- Header bar (full-width, couleur primaire) ----

def _header(slide, title, subtitle=None, height_in=1.15, brand=None):
    primary = (brand or {}).get("primary", NAVY)
    font = (brand or _CURRENT_BRAND or {}).get("font", FONT)
    h = Inches(height_in)
    _rect(slide, Emu(0), Emu(0), SW, h, fill=primary)
    tb = _tb(slide, Inches(0.4), Inches(0.12), SW - Inches(0.8), h - Inches(0.25))
    _clear(tb.text_frame)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.name = font
    r.font.color.rgb = WHITE
    if subtitle:
        _run(tb.text_frame, subtitle, size=9, color=RGBColor(0xAA, 0xCC, 0xFF))


# ---- Section label ----

def _section_label(slide, left, top, width, height_in, label):
    h = Inches(height_in)
    _rect(slide, left, top, width, h, fill=LIGHT_BLUE, no_line=False, line_color=BORDER, line_pt=0.3)
    tb = _tb(slide, left + Inches(0.12), top + Inches(0.04), width - Inches(0.2), h)
    _clear(tb.text_frame)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = label.upper()
    r.font.size = Pt(7)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = BLUE


# ---- Slide footer ----

def _footer(slide, text="MARCEL · Confidentiel"):
    tb = _tb(slide, Inches(0.4), SH - Inches(0.3), SW - Inches(1.8), Inches(0.25))
    _clear(tb.text_frame)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(6.5)
    r.font.name = _CURRENT_BRAND.get("font", FONT) if _CURRENT_BRAND else FONT
    r.font.color.rgb = LIGHT
    logo_b64 = _CURRENT_BRAND.get("logo_base64") if _CURRENT_BRAND else None
    if logo_b64:
        try:
            logo_data = base64.b64decode(logo_b64)
            logo_buf = io.BytesIO(logo_data)
            logo_w = Inches(0.7)
            slide.shapes.add_picture(
                logo_buf, SW - logo_w - Inches(0.15), SH - Inches(0.35), width=logo_w
            )
        except Exception:
            pass


# ---- Table helper ----

def _styled_table(slide, headers, rows, left, top, width, col_widths_in, row_height_in=0.42,
                  header_bg=NAVY, header_fg=WHITE):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    h_total = Inches(row_height_in) * n_rows
    tf_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, h_total)
    tf = tf_shape.table
    for i, w in enumerate(col_widths_in):
        tf.columns[i].width = Inches(w)
    for row in tf.rows:
        row.height = Inches(row_height_in)
    for col, h in enumerate(headers):
        cell = tf.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.02)
        cell.text_frame.word_wrap = True
        cell.text_frame.paragraphs[0].clear()
        r = cell.text_frame.paragraphs[0].add_run()
        r.text = h
        r.font.size = Pt(8)
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = header_fg
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for row_idx, row_data in enumerate(rows):
        row_num = row_idx + 1
        row_bg = BG if row_idx % 2 == 0 else WHITE
        for col, (val, opts) in enumerate(row_data):
            cell = tf.cell(row_num, col)
            cell.fill.solid()
            bg = opts.get("bg", row_bg)
            cell.fill.fore_color.rgb = bg
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.text_frame.word_wrap = True
            cell.text_frame.paragraphs[0].clear()
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = str(val)
            run.font.size = Pt(opts.get("size", 8))
            run.font.bold = opts.get("bold", False)
            run.font.name = FONT
            run.font.color.rgb = opts.get("color", DARK)
            cell.text_frame.paragraphs[0].alignment = opts.get("align", PP_ALIGN.LEFT)
    return tf


# ---- 2-column key/value table helper ----

def _kv_table(slide, items, left, top, width, label_ratio=0.45, row_h=0.36):
    """items = list of (label, value, value_color, bold, row_bg)"""
    n = len(items)
    height = Inches(row_h * n)
    ts = slide.shapes.add_table(n, 2, left, top, width, height)
    tbl = ts.table
    lw = int(width * label_ratio)
    tbl.columns[0].width = lw
    tbl.columns[1].width = width - lw
    for i, (lbl, val, fg, bold, bg) in enumerate(items):
        tbl.rows[i].height = Inches(row_h)
        for col_idx, (text, clr, is_bold) in enumerate([(lbl, MID, False), (val, fg, bold)]):
            cell = tbl.cell(i, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.margin_left = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.text_frame.word_wrap = True
            p = cell.text_frame.paragraphs[0]
            p.clear()
            r = p.add_run()
            r.text = text
            r.font.size = Pt(8.5)
            r.font.bold = is_bold
            r.font.name = FONT
            r.font.color.rgb = clr
    return ts, Inches(row_h * n)


# ---- S2-04 — Slide Roadmap matplotlib ----

