"""Backend tests for the Engagement/Gate criteria system, readiness, gate blocking, and PPTX export."""
import os
import io
import pytest
import requests

BASE_URL = os.environ.get("REACT_APP_BACKEND_URL").rstrip("/")
API = f"{BASE_URL}/api"

ADMIN = {"email": "admin@altair.fr", "password": "Admin2026!"}


# ─── Fixtures ─────────────────────────────────────────────────────────────────
@pytest.fixture(scope="module")
def token():
    r = requests.post(f"{API}/auth/login", json=ADMIN, timeout=15)
    assert r.status_code == 200, r.text
    return r.json()["access_token"]


@pytest.fixture(scope="module")
def h(token):
    return {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}


@pytest.fixture(scope="module")
def cadrage_project(h):
    """Find a project in cadrage phase, or the first project."""
    r = requests.get(f"{API}/projects", headers=h, timeout=15)
    assert r.status_code == 200
    projs = r.json()
    cad = [p for p in projs if (p.get("lifecycle_phase") or "cadrage") == "cadrage"]
    return (cad or projs)[0]


# ─── Auth guard ───────────────────────────────────────────────────────────────
def test_criteria_requires_auth():
    r = requests.get(f"{API}/engagement/criteria/cadrage", timeout=10)
    assert r.status_code in (401, 403)


def test_pptx_requires_auth(cadrage_project):
    r = requests.get(f"{API}/exports/engagement/{cadrage_project['project_id']}.pptx", timeout=15)
    assert r.status_code in (401, 403)


# ─── Criteria list & seed ─────────────────────────────────────────────────────
def test_criteria_cadrage_21(h):
    r = requests.get(f"{API}/engagement/criteria/cadrage", headers=h, timeout=15)
    assert r.status_code == 200
    data = r.json()
    assert len(data) == 21, f"cadrage expected 21 criteria, got {len(data)}"
    types = {c["type"] for c in data}
    assert types == {"auto", "attested"}
    assert any(c["mandatory"] for c in data)


def test_criteria_conception_29(h):
    r = requests.get(f"{API}/engagement/criteria/conception", headers=h, timeout=15)
    assert r.status_code == 200
    assert len(r.json()) == 29


def test_criteria_light_phases(h):
    for phase in ("realisation", "recette"):
        r = requests.get(f"{API}/engagement/criteria/{phase}", headers=h, timeout=15)
        assert r.status_code == 200
        n = len(r.json())
        assert 3 <= n <= 10, f"{phase}: {n}"
    r = requests.get(f"{API}/engagement/criteria/deploiement", headers=h, timeout=15)
    assert r.status_code == 200
    # deploiement = 1 + light
    assert len(r.json()) >= 5


# ─── CRUD critères ────────────────────────────────────────────────────────────
def test_criteria_crud_custom(h):
    # Create custom
    r = requests.post(f"{API}/engagement/criteria", headers=h,
                      json={"from_phase": "cadrage", "label": "TEST Custom criterion", "mandatory": False}, timeout=15)
    assert r.status_code == 200, r.text
    crit = r.json()
    assert crit["custom"] is True
    assert crit["type"] == "attested"
    cid = crit["criterion_id"]

    # Update label
    r2 = requests.patch(f"{API}/engagement/criteria/{cid}", headers=h,
                        json={"label": "TEST Custom updated", "mandatory": True}, timeout=15)
    assert r2.status_code == 200
    assert r2.json()["label"] == "TEST Custom updated"
    assert r2.json()["mandatory"] is True

    # Delete custom
    r3 = requests.delete(f"{API}/engagement/criteria/{cid}", headers=h, timeout=15)
    assert r3.status_code == 204


def test_cannot_delete_default_criterion(h):
    r = requests.get(f"{API}/engagement/criteria/cadrage", headers=h, timeout=15)
    default_c = next(c for c in r.json() if not c.get("custom"))
    r2 = requests.delete(f"{API}/engagement/criteria/{default_c['criterion_id']}", headers=h, timeout=15)
    assert r2.status_code == 400


# ─── Readiness ────────────────────────────────────────────────────────────────
def test_readiness_structure(h, cadrage_project):
    pid = cadrage_project["project_id"]
    r = requests.get(f"{API}/projects/{pid}/engagement/readiness", headers=h, timeout=15)
    assert r.status_code == 200, r.text
    d = r.json()
    for k in ("from_phase", "items", "score_pct", "ready", "mandatory_missing"):
        assert k in d
    assert isinstance(d["items"], list) and len(d["items"]) > 0
    it = d["items"][0]
    for k in ("ok", "mandatory", "type", "attestation"):
        assert k in it
    assert 0 <= d["score_pct"] <= 100


# ─── Attest ──────────────────────────────────────────────────────────────────
def test_attest_raises_score_and_cleanup(h, cadrage_project):
    pid = cadrage_project["project_id"]
    r0 = requests.get(f"{API}/projects/{pid}/engagement/readiness", headers=h, timeout=15).json()
    # Find an attested criterion currently not ok
    tgt = next((it for it in r0["items"] if it["type"] == "attested" and not it["ok"]), None)
    if not tgt:
        pytest.skip("No unattested criterion available")
    score_before, miss_before = r0["score_pct"], len(r0["mandatory_missing"])

    # Attest checked=True
    ra = requests.post(f"{API}/projects/{pid}/engagement/attest", headers=h,
                       json={"criterion_id": tgt["criterion_id"], "checked": True,
                             "justification": "TEST justification"}, timeout=15)
    assert ra.status_code == 200

    r1 = requests.get(f"{API}/projects/{pid}/engagement/readiness", headers=h, timeout=15).json()
    assert r1["score_pct"] >= score_before
    if tgt["mandatory"]:
        assert len(r1["mandatory_missing"]) <= miss_before

    # not_applicable path
    rn = requests.post(f"{API}/projects/{pid}/engagement/attest", headers=h,
                       json={"criterion_id": tgt["criterion_id"], "checked": False,
                             "not_applicable": True, "justification": "TEST NA"}, timeout=15)
    assert rn.status_code == 200

    # Cleanup: uncheck
    requests.post(f"{API}/projects/{pid}/engagement/attest", headers=h,
                  json={"criterion_id": tgt["criterion_id"], "checked": False,
                        "not_applicable": False, "justification": ""}, timeout=15)


# ─── Gate blocking (422) ─────────────────────────────────────────────────────
def test_gate_blocked_without_override(h, cadrage_project):
    pid = cadrage_project["project_id"]
    rdy = requests.get(f"{API}/projects/{pid}/engagement/readiness", headers=h, timeout=15).json()
    if not rdy["mandatory_missing"]:
        pytest.skip("Project already ready — cannot test 422")

    # Cancel any open gate first
    lc = requests.get(f"{API}/projects/{pid}/lifecycle", headers=h, timeout=15).json()
    for g in lc.get("gates", []):
        if g["status"] in ("en_validation", "pret"):
            requests.delete(f"{API}/lifecycle/gates/{g['gate_id']}", headers=h, timeout=15)

    r = requests.post(f"{API}/projects/{pid}/lifecycle/gates", headers=h,
                      json={"target_date": "2026-12-31"}, timeout=15)
    assert r.status_code == 422, f"expected 422, got {r.status_code}: {r.text}"
    d = r.json().get("detail", r.json())
    assert isinstance(d, dict)
    assert "score_pct" in d and "mandatory_missing" in d


def test_gate_created_with_override_and_cleanup(h, cadrage_project):
    pid = cadrage_project["project_id"]
    rdy = requests.get(f"{API}/projects/{pid}/engagement/readiness", headers=h, timeout=15).json()
    if not rdy["mandatory_missing"]:
        pytest.skip("Project already ready — override not exercised")

    # Ensure no existing open gate
    lc = requests.get(f"{API}/projects/{pid}/lifecycle", headers=h, timeout=15).json()
    for g in lc.get("gates", []):
        if g["status"] in ("en_validation", "pret"):
            requests.delete(f"{API}/lifecycle/gates/{g['gate_id']}", headers=h, timeout=15)

    r = requests.post(f"{API}/projects/{pid}/lifecycle/gates", headers=h,
                      json={"target_date": "2026-12-31", "readiness_override": True}, timeout=15)
    assert r.status_code in (200, 201), r.text
    gate = r.json()
    assert gate["readiness_override"] is True
    assert "readiness_score" in gate
    gate_id = gate["gate_id"]

    # Cleanup: cancel the gate
    dr = requests.delete(f"{API}/lifecycle/gates/{gate_id}", headers=h, timeout=15)
    assert dr.status_code in (200, 204)


# ─── Project fields persistence ──────────────────────────────────────────────
def test_project_new_fields_persistence(h, cadrage_project):
    pid = cadrage_project["project_id"]

    # Snapshot original values to restore
    orig = requests.get(f"{API}/projects/{pid}", headers=h, timeout=15).json()

    payload = {
        "scope_in": "TEST scope in",
        "scope_out": "TEST scope out",
        "nfr": "TEST NFR content",
        "impacted_entities": ["Entity A", "Entity B"],
        "governance_roles": [{"role": "Sponsor", "name": "TEST Sponsor"},
                             {"role": "Chef de projet", "name": "TEST CP"}],
        "build_to_run": "TEST build-to-run",
        "budget_breakdown": [{"entity": "Entity A", "capex": 100, "opex": 20}],
    }
    r = requests.put(f"{API}/projects/{pid}", headers=h, json=payload, timeout=15)
    assert r.status_code == 200, r.text

    got = requests.get(f"{API}/projects/{pid}", headers=h, timeout=15).json()
    for k, v in payload.items():
        assert got.get(k) == v, f"field {k} not persisted: {got.get(k)!r}"

    # Restore
    restore = {k: orig.get(k) for k in payload}
    requests.put(f"{API}/projects/{pid}", headers=h, json=restore, timeout=15)


# ─── PPTX export ─────────────────────────────────────────────────────────────
def test_engagement_pptx_download(h, cadrage_project):
    pid = cadrage_project["project_id"]
    r = requests.get(f"{API}/exports/engagement/{pid}.pptx", headers=h, timeout=30)
    assert r.status_code == 200, r.text[:500]
    assert r.content[:2] == b"PK", "Not a valid PPTX (zip) payload"
    # Verify slide count is around 10
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(r.content))
        assert len(prs.slides) >= 8, f"Only {len(prs.slides)} slides"
    except ImportError:
        pass
